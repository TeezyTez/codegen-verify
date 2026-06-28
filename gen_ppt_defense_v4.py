# -*- coding: utf-8 -*-
"""开题答辩PPT v4.1 — 去AI味"""
import sys; sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation; from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor; from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE; from pptx.oxml.ns import qn
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height

DARK = RGBColor(0x1B,0x3A,0x5C); MID = RGBColor(0x2C,0x5F,0x8A); BLUE = RGBColor(0x3A,0x7C,0xBD)
ORANGE = RGBColor(0xE8,0x6C,0x00); WHITE = RGBColor(0xFF,0xFF,0xFF); BLACK = RGBColor(0x33,0x33,0x33)
GRAY = RGBColor(0x88,0x88,0x88); LGRAY = RGBColor(0xF2,0xF3,0xF5); RED = RGBColor(0xCC,0x33,0x33)
GREEN = RGBColor(0x2E,0x7D,0x32); TEAL = RGBColor(0x00,0x7A,0x6B)
TF='SimHei'; BF='SimSun'; EF='Calibri'

def blank(): return prs.slides.add_slide(prs.slide_layouts[6])
def fs(s,c=DARK):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H); sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
def rc(s,l,t,w,h,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h); sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
def rd(s,l,t,w,h,c):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h); sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
def tx(s,l,t,w,h,txt,fn=BF,fs_n=16,b=False,c=BLACK,al=PP_ALIGN.LEFT,ls=1.2):
    bx=s.shapes.add_textbox(l,t,w,h); bx.text_frame.word_wrap=True
    p=bx.text_frame.paragraphs[0]; p.text=txt; p.font.name=fn; p.font.size=Pt(fs_n); p.font.bold=b; p.font.color.rgb=c; p.alignment=al
    pPr=p._p.get_or_add_pPr(); ln=pPr.makeelement(qn('a:lnSpc'),{})
    sp=ln.makeelement(qn('a:spcPct'),{'val':str(int(ls*100000))}); ln.append(sp); pPr.append(ln)
def ov(s,l,t,sz,num,label,cl=DARK):
    c=s.shapes.add_shape(MSO_SHAPE.OVAL,l,t,sz,sz); c.fill.solid(); c.fill.fore_color.rgb=cl; c.line.fill.background()
    tf=c.text_frame; tf.paragraphs[0].text=num; tf.paragraphs[0].font.name=EF; tf.paragraphs[0].font.size=Pt(28); tf.paragraphs[0].font.bold=True
    tf.paragraphs[0].font.color.rgb=WHITE; tf.paragraphs[0].alignment=PP_ALIGN.CENTER; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    tx(s,l+sz+Inches(0.2),t+Inches(0.15),Inches(3),Inches(0.4),label,TF,18,False,DARK)

# ====== S1 封面 ======
s=blank(); fs(s,DARK)
rc(s,0,Inches(0.15),W,Inches(0.05),ORANGE)
tx(s,Inches(1),Inches(0.8),Inches(11),Inches(0.5),"华东师范大学  \u00b7  软件工程",TF,18,False,RGBColor(0x8A,0xB4,0xF8),PP_ALIGN.CENTER)
tx(s,Inches(1),Inches(2.0),Inches(11),Inches(1.5),"基于形式化规约的 LLM\n代码生成与自修复方法研究",TF,42,True,WHITE,PP_ALIGN.CENTER,1.3)
rc(s,Inches(5.5),Inches(3.8),Inches(2.3),Inches(0.04),ORANGE)
mt=[("硕士论文开题答辩",BF,18,False,RGBColor(0xCC,0xD5,0xE8),PP_ALIGN.CENTER),("",BF,14,False,WHITE,PP_ALIGN.CENTER),("Tez  \u00b7  导师：XXX 教授",BF,18,False,WHITE,PP_ALIGN.CENTER),("2026年6月",BF,16,False,RGBColor(0x8A,0xB4,0xF8),PP_ALIGN.CENTER)]
bx=s.shapes.add_textbox(Inches(3),Inches(4.2),Inches(7),Inches(2.0)); bx.text_frame.word_wrap=True
for i,(t,fn,fs2,b,c,al) in enumerate(mt):
    p=bx.text_frame.paragraphs[0] if i==0 else bx.text_frame.add_paragraph()
    p.text=t; p.font.name=fn; p.font.size=Pt(fs2); p.font.bold=b; p.font.color.rgb=c; p.alignment=al
rc(s,0,H-Inches(0.15),W,Inches(0.05),ORANGE)
print("S1 OK")

# ====== S2 目录 ======
s=blank()
tx(s,Inches(0.8),Inches(0.3),Inches(3),Inches(0.6),"目录",TF,36,True,DARK)
rc(s,Inches(0.8),Inches(0.9),Inches(0.8),Inches(0.04),ORANGE)
its=[("01","研究背景"),("02","高pass@1不等于正确性"),("03","方案与架构"),("04","关键技术"),("05","实验与进展"),("06","创新点与时间线")]
colors=[DARK,MID,BLUE,ORANGE,MID,DARK]
for i,(n,t) in enumerate(its):
    row,col=i//3,i%3; x=Inches(1.5)+col*Inches(3.5); y=Inches(1.8)+row*Inches(2.5)
    shp=s.shapes.add_shape(MSO_SHAPE.OVAL if i!=3 else MSO_SHAPE.ROUNDED_RECTANGLE,x,y,Inches(0.9),Inches(0.9))
    shp.fill.solid(); shp.fill.fore_color.rgb=colors[i]; shp.line.fill.background()
    tf=shp.text_frame; tf.paragraphs[0].text=n; tf.paragraphs[0].font.name=EF; tf.paragraphs[0].font.size=Pt(24)
    tf.paragraphs[0].font.bold=True; tf.paragraphs[0].font.color.rgb=WHITE; tf.paragraphs[0].alignment=PP_ALIGN.CENTER; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    tx(s,x+Inches(1.1),y+Inches(0.25),Inches(2.5),Inches(0.5),t,TF,20,True,DARK)
print("S2 OK")

# ====== S3 背景 ======
s=blank()
rc(s,0,0,W,Inches(0.06),DARK)
tx(s,Inches(1),Inches(0.3),Inches(11),Inches(0.6),"故事从这开始",TF,30,True,DARK)
rc(s,Inches(1),Inches(0.9),Inches(0.8),Inches(0.03),ORANGE)
ov(s,Inches(0.8),Inches(1.5),Inches(1.0),"28.8%","Codex 2021",DARK)
ov(s,Inches(0.8),Inches(3.0),Inches(1.0),"94.5%","ReflexiCoder 2026",MID)
ov(s,Inches(0.8),Inches(4.5),Inches(1.0),"100\u4e07+","\u7ec4\u7ec7\u7528 GitHub Copilot",BLUE)
tx(s,Inches(5.5),Inches(1.5),Inches(7),Inches(0.5),"LLM 写代码越来越强，但\u2026\u2026",TF,22,True,RED)
tx(s,Inches(5.5),Inches(2.1),Inches(7),Inches(1.5),
   "测试覆盖是最大的谎言：\n"
   "HumanEval 每题只有 7.8 个测试用例。\n"
   "扩到 774 个 \u2192 pass@1 从 85.4% 掉到 33.5%。\n\n"
   "模型能\u2018写\u2019代码，但本质上不\u2018理解\u2019它在写什么。",BF,16,False,BLACK,ls=1.5)
rd(s,Inches(5.5),Inches(4.2),Inches(7),Inches(1.2),LGRAY)
tx(s,Inches(5.8),Inches(4.3),Inches(6.5),Inches(1.0),
   "\U0001f4a1 核心问题：生成能力越强，我们越需要一个办法来回答\n"
   "\u2014\u2014 这一段具体代码，到底是不是对的？",TF,17,True,DARK,ls=1.3)
print("S3 OK")

# ====== S4 三个危机 ======
s=blank()
rc(s,0,0,W,Inches(0.06),DARK)
tx(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.6),"代码自修复的三个死穴",TF,30,True,DARK)
rc(s,Inches(0.8),Inches(0.9),Inches(0.8),Inches(0.03),ORANGE)

crises_data = [
    ("\U0001f50d 自审查","双向失效",DARK,["漏检率 31.7%","误拒率飙到 73%","越详细提示越糟糕","\u2192 Reddy 2026 / Jin 2026"]),
    ("\U0001f9ea 反馈质量","先天不足",MID,["人类反馈 > GPT-4 乘以1.58","断言错误修复率仅 45%","指令微调反而削了修复力","\u2192 Olausson 2023 / Arimbur 2026"]),
    ("\U0001f6e1 验证手段","没有独立裁判",BLUE,["测试覆盖不全","LLM 内省共享同样盲区","缺少数学级的验证","\u2192 这就是我的切入点"]),
]
for i,(icon,t,cl,items) in enumerate(crises_data):
    x=Inches(0.6)+i*Inches(4.2); w=Inches(3.9)
    rc(s,x,Inches(1.3),w,Inches(0.06),cl)
    rc(s,x,Inches(1.36),w,Inches(5.0),LGRAY)
    tx(s,x+Inches(0.2),Inches(1.5),w-Inches(0.4),Inches(0.5),f"{icon}  {t}",TF,20,True,cl)
    rc(s,x+Inches(0.3),Inches(2.1),w-Inches(0.6),Inches(0.02),ORANGE)
    for j,item in enumerate(items):
        tx(s,x+Inches(0.3),Inches(2.4+j*0.7),w-Inches(0.6),Inches(0.6),item,BF,16,False,BLACK if j<3 else ORANGE,ls=1.2)
print("S4 OK")

# ====== S5 高pass@1 ======
s=blank(); fs(s,DARK)
rc(s,0,Inches(0.15),W,Inches(0.05),ORANGE)
tx(s,Inches(1),Inches(0.6),Inches(11),Inches(0.6),"那 94.5% 的 HumanEval 怎么说？",TF,34,True,WHITE,PP_ALIGN.CENTER)
rc(s,Inches(4.5),Inches(1.3),Inches(4.3),Inches(0.04),ORANGE)

quotes=[
    ("第一","测试覆盖是假象","每题平均 7.8 个测试 \u2192 774 个后正确率从 85% 掉到 33%",DARK),
    ("第二","失败的那 5% 才是关键","断言错误修复率 45%，自审漏检 31.7%\u2014\u2014高正确率治不了这个",MID),
    ("第三","高 pass@1 是我的前提，不是对手","LLM 够强才能自动写规约。我利用它，不是否定它",BLUE),
]
for i,(n,t,d,cl) in enumerate(quotes):
    y=Inches(1.8+i*1.7)
    rc(s,Inches(1),y,Inches(1.2),Inches(1.3),cl)
    tx(s,Inches(1),y+Inches(0.4),Inches(1.2),Inches(0.5),n,TF,20,True,WHITE,PP_ALIGN.CENTER)
    rc(s,Inches(2.4),y,Inches(10.3),Inches(1.3),RGBColor(0x2A,0x4A,0x6E))
    tx(s,Inches(2.7),y+Inches(0.1),Inches(9.7),Inches(0.4),t,TF,20,True,WHITE)
    tx(s,Inches(2.7),y+Inches(0.55),Inches(9.7),Inches(0.6),d,BF,16,False,RGBColor(0xBB,0xCC,0xDD))

rc(s,Inches(1),Inches(5.8),Inches(11.3),Inches(1.0),RGBColor(0x15,0x2D,0x4A))
tx(s,Inches(1.2),Inches(5.85),Inches(10.9),Inches(0.8),
   "高 pass@1 回答的是「模型平均能写对多少」  形式化验证回答的是「这段代码确实正确吗」\n"
   "二者正交。安全关键场景下，后者不可或缺。",TF,17,True,RGBColor(0x8A,0xB4,0xF8),PP_ALIGN.CENTER,1.4)
print("S5 OK")

# ====== S6 形式化方案 ======
s=blank()
rc(s,0,0,W,Inches(0.06),DARK)
tx(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.6),"我的解法：引入一个不犯糊涂的裁判",TF,28,True,DARK)
rc(s,Inches(0.8),Inches(0.9),Inches(0.8),Inches(0.03),ORANGE)
tx(s,Inches(0.8),Inches(1.5),Inches(4),Inches(0.4),"自修复方案的三种反馈对比",TF,18,True,DARK)

# Three comparison blocks - each: title, icon, items(list), color
comps_data = [
    ("测试执行","\u274c",["只告诉你「没过」","覆盖不到的地方就是盲区"],GRAY),
    ("LLM 自审","\u274c",["漏检 31.7% 的 bug","还会把对的当成错的"],GRAY),
    ("Dafny 验证器","\u2705",["\u2713 精确到行号和列","\u2713 告诉你违反了什么规约","\u2713 数学级别的正确性保证"],DARK),
]
for i,(t,icon,items,cl) in enumerate(comps_data):
    x=Inches(0.8+i*4.1)
    rc(s,x,Inches(2.1),Inches(3.8),Inches(3.0),LGRAY)
    tx(s,x+Inches(0.2),Inches(2.2),Inches(3.4),Inches(0.4),f"{icon}  {t}",TF,18,True,cl)
    rc(s,x+Inches(0.3),Inches(2.7),Inches(3.2),Inches(0.02),ORANGE)
    for j,item in enumerate(items):
        tx(s,x+Inches(0.2),Inches(2.9+j*0.6),Inches(3.4),Inches(0.5),item,BF,15,False,BLACK,ls=1.4)

rd(s,Inches(0.8),Inches(5.5),Inches(11.7),Inches(1.0),DARK)
tx(s,Inches(1.2),Inches(5.6),Inches(11.1),Inches(0.8),
   "关键思路：规约（spec）是「什么算对」的精确回答  \u2192 验证器是「对不对」的独立裁判  \u2192\n"
   "Diagnose Agent 把裁判的话翻给 LLM 听懂  \u2192 LLM 自己改",TF,17,True,WHITE,PP_ALIGN.CENTER,1.3)
print("S6 OK")

# ====== S7 研究目标 ======
s=blank()
rc(s,0,0,W,Inches(0.06),DARK)
tx(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.6),"研究目标：四件事",TF,30,True,DARK)
rc(s,Inches(0.8),Inches(0.9),Inches(0.8),Inches(0.03),ORANGE)
tx(s,Inches(0.8),Inches(1.3),Inches(11),Inches(0.5),
   "从自然语言 \u2192 Dafny 规约 \u2192 规约约束下生成代码 \u2192 验证器反馈驱动修复 \u2192 可证明正确的代码",TF,16,False,GRAY,ls=1.2)

goals=[("一","规约自动生成","NL 描述 \u2192 Dafny requires/ensures\n自带 Dafny resolve 合法性校验",DARK),
       ("二","规约感知的代码生成","规约作为第一优先级约束\n嵌入代码生成 Prompt",MID),
       ("三","验证反馈结构化翻译","Dafny 原始错误\n\u2192 分类 / 定位 / 修复引导",BLUE),
       ("四","退化感知的迭代修复","Coordinator 控退化\nMemory Agent 记经验",ORANGE)]
for i,(n,t,d,cl) in enumerate(goals):
    x=Inches(0.5+i*3.2)
    rc(s,x,Inches(2.0),Inches(3.0),Inches(0.6),cl)
    tx(s,x,Inches(2.05),Inches(3.0),Inches(0.5),f"内容{n}  {t}",TF,16,True,WHITE,PP_ALIGN.CENTER)
    rc(s,x,Inches(2.6),Inches(3.0),Inches(2.5),LGRAY)
    tx(s,x+Inches(0.2),Inches(2.8),Inches(2.6),Inches(2.0),d,BF,15,False,BLACK,ls=1.5)
    rc(s,x+Inches(0.2),Inches(4.2),Inches(2.6),Inches(0.7),WHITE)
    tx(s,x+Inches(0.3),Inches(4.25),Inches(2.4),Inches(0.6),
       "\u2753 关键：如何确保规约既完整又正确？",TF,11,True,ORANGE,PP_ALIGN.CENTER,1.2)
print("S7 OK")

# ====== S8 流程图 ======
s=blank()
rc(s,0,0,W,Inches(0.06),DARK)
tx(s,Inches(0.8),Inches(0.2),Inches(11),Inches(0.5),"流程图",TF,28,True,DARK)
for lbl,x,w,cl in [("Phase 1: 规约生成",0.5,4.0,DARK),("Phase 2: 代码生成",4.7,3.5,MID),("Phase 3: 验证\u2192修复循环",8.5,4.5,ORANGE)]:
    rc(s,Inches(x),Inches(1.2),Inches(w),Inches(0.35),cl)
    tx(s,Inches(x),Inches(1.22),Inches(w),Inches(0.3),lbl,TF,12,True,WHITE,PP_ALIGN.CENTER)

f_items=[("输入: 问题描述",MID),("Spec Agent\n生成规约",ORANGE),("Dafny\n检查规约",MID),
         ("Code Agent\n生成代码",ORANGE),("Dafny\n验证代码",GREEN),("END \u2705",GREEN)]
for i,(t,cl) in enumerate(f_items):
    x=Inches(0.5+i*2.5)
    rd(s,x,Inches(1.9),Inches(2.0),Inches(0.55),cl)
    tx(s,x,Inches(1.92),Inches(2.0),Inches(0.55),t,TF,13,True,WHITE,PP_ALIGN.CENTER)
    s.shapes[-1].text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE

rc(s,Inches(0.5),Inches(2.8),Inches(11.7),Inches(0.02),RED)
tx(s,Inches(9.5),Inches(2.6),Inches(3),Inches(0.3),"\u274c 进入修复循环 \u2192",BF,14,True,RED,PP_ALIGN.RIGHT)
r_items=[("Diagnose Agent\n分析错误",MID),("Repair Agent\n执行修复",ORANGE),("Coordinator\n控退化, 调策略",DARK),("Memory Agent\n检索历史经验",BLUE)]
for i,(t,cl) in enumerate(r_items):
    x=Inches(0.5+i*3.2)
    rd(s,x,Inches(3.2),Inches(2.8),Inches(0.5),cl)
    tx(s,x,Inches(3.22),Inches(2.8),Inches(0.5),t,TF,12,True,WHITE,PP_ALIGN.CENTER)
    s.shapes[-1].text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
tx(s,Inches(9.5),Inches(3.9),Inches(3.5),Inches(0.3),"\u2190 修复完回去验证 \u2192",BF,12,True,GRAY,PP_ALIGN.CENTER)
rc(s,Inches(9.5),Inches(4.2),Inches(3.5),Inches(0.02),GRAY)
rd(s,Inches(0.5),Inches(4.6),Inches(11.7),Inches(0.5),LGRAY)
tx(s,Inches(0.8),Inches(4.65),Inches(11.2),Inches(0.4),
   "最多 3 轮。Coordinator 如果发现越修越差 \u2192 自动回退到上一个最好版本",BF,14,False,BLACK,PP_ALIGN.CENTER)
tx(s,Inches(0.5),Inches(6.5),Inches(12),Inches(0.4),
   "技术栈：LangGraph | DeepSeek / GPT | Dafny CLI | Python 3.11+",BF,12,False,GRAY,PP_ALIGN.CENTER)
print("S8 OK")

# ====== S9 关键技术 ======
s=blank()
rc(s,0,0,W,Inches(0.06),DARK)
tx(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.6),"核心技术：翻译验证器的话",TF,28,True,DARK)
rc(s,Inches(0.8),Inches(0.9),Inches(0.8),Inches(0.03),ORANGE)
tx(s,Inches(0.8),Inches(1.4),Inches(5),Inches(0.4),"Dafny 说的是\u2026\u2026",TF,18,True,RED)
rd(s,Inches(0.8),Inches(1.9),Inches(5.5),Inches(2.5),RGBColor(0xFD,0xED,0xED))
tx(s,Inches(1.0),Inches(2.0),Inches(5.1),Inches(2.3),
   "\u300ctest.dfy(10,4): Error: A postcondition\nmight not hold on this return path.\u300d\n\n"
   "人类看得懂吗？LLM 看得懂吗？\n\u2014\u2014 这就是问题所在。",BF,15,False,RED,ls=1.4)
tx(s,Inches(7.0),Inches(1.4),Inches(5.5),Inches(0.4),"Diagnose Agent 翻译成\u2026\u2026",TF,18,True,GREEN)
rd(s,Inches(7.0),Inches(1.9),Inches(5.5),Inches(2.5),RGBColor(0xE8,0xF5,0xE9))
tx(s,Inches(7.2),Inches(2.0),Inches(5.1),Inches(2.3),
   "\u300c第 10 行第 4 列 \u2192 后置条件没满足\n关联规约：ensures result >= x\n修复建议：检查该分支上 result 的计算\u300d\n\n"
   "\u2192 LLM 直接就能用了。",BF,15,False,RGBColor(0x2E,0x7D,0x32),ls=1.4)
arr=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(6.4),Inches(2.5),Inches(0.5),Inches(0.3))
arr.fill.solid(); arr.fill.fore_color.rgb=ORANGE; arr.line.fill.background()
tx(s,Inches(6.2),Inches(2.85),Inches(0.8),Inches(0.3),"\u7ffb\u8bd1",BF,10,True,ORANGE,PP_ALIGN.CENTER)
rc(s,Inches(0.8),Inches(4.8),Inches(11.7),Inches(1.5),LGRAY)
tx(s,Inches(1.0),Inches(4.9),Inches(11.3),Inches(0.4),"三类最常见错误 + 修复策略",TF,18,True,DARK)
errs=[("postcondition","检查函数逻辑是否正确","#CC3333"),("invariant","调整循环不变量或循环体","#2C5F8A"),("syntax / type","直接修正语法或类型标注","#336633")]
for i,(et,fix,hc) in enumerate(errs):
    x=Inches(1.0+i*3.8); r=int(hc[1:3],16); g=int(hc[3:5],16); b2=int(hc[5:7],16)
    rc(s,x,Inches(5.5),Inches(3.4),Inches(0.5),RGBColor(r,g,b2))
    tx(s,x,Inches(5.52),Inches(3.4),Inches(0.45),et,TF,15,True,WHITE,PP_ALIGN.CENTER)
    tx(s,x,Inches(6.1),Inches(3.4),Inches(0.3),fix,BF,13,False,BLACK,PP_ALIGN.CENTER)
print("S9 OK")

# ====== S10 实验方案 ======
s=blank()
rc(s,0,0,W,Inches(0.06),DARK)
tx(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.6),"怎么验证我做得对不对",TF,28,True,DARK)
rc(s,Inches(0.8),Inches(0.9),Inches(0.8),Inches(0.03),ORANGE)
rd(s,Inches(0.8),Inches(1.3),Inches(11.7),Inches(0.6),LGRAY)
tx(s,Inches(1.0),Inches(1.35),Inches(11.3),Inches(0.5),
   "数据集：HumanEval（164题，主实验）  |  MBPP（974题，泛化性验证）  |  Proving the Coding Interview（27题，质量参照）",BF,15,False,DARK,PP_ALIGN.CENTER)
tx(s,Inches(0.8),Inches(2.2),Inches(5),Inches(0.4),"要跟谁比",TF,20,True,DARK)
bls=["Direct Gen","Self-Debug","Reflexion","ReflexiCoder","单 Agent"]
for i,b in enumerate(bls):
    x=Inches(0.8+i*2.5)
    rd(s,x,Inches(2.7),Inches(2.3),Inches(0.55),DARK if i%2==0 else MID)
    tx(s,x,Inches(2.72),Inches(2.3),Inches(0.5),b,TF,15,True,WHITE,PP_ALIGN.CENTER)
tx(s,Inches(0.8),Inches(3.6),Inches(5),Inches(0.4),"看什么指标",TF,20,True,DARK)
rd(s,Inches(0.8),Inches(4.1),Inches(11.7),Inches(0.55),LGRAY)
tx(s,Inches(1.0),Inches(4.12),Inches(11.3),Inches(0.5),
   "Verif@1（首轮通过率） | Verif@k（k轮通过率） | SpecAcc（规约正确率） | AvgRounds（平均轮次） | DegradationRate（退化率）",BF,14,False,BLACK,PP_ALIGN.CENTER)
tx(s,Inches(0.8),Inches(5.0),Inches(5),Inches(0.4),"消融实验拆开看",TF,20,True,DARK)
abls=[("A1 规约有无","规约对首轮通过率的贡献"),("A2 反馈结构","结构化翻译 vs 原始输出"),("A3 Agent分离","单 Agent vs 多 Agent"),("A4 轮次上限","1/3/5轮效果对比"),("A5 模型对比","DeepSeek/GPT/Qwen3")]
for i,(n,d) in enumerate(abls):
    x=Inches(0.8+i*2.5)
    rd(s,x,Inches(5.5),Inches(2.3),Inches(0.85),LGRAY)
    tx(s,x+Inches(0.1),Inches(5.55),Inches(2.1),Inches(0.3),n,TF,15,True,ORANGE)
    tx(s,x+Inches(0.1),Inches(5.85),Inches(2.1),Inches(0.4),d,BF,13,False,BLACK)
print("S10 OK")

# ====== S11 当前进展 ======
s=blank()
rc(s,0,0,W,Inches(0.06),DARK)
tx(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.6),"已经做了什么",TF,30,True,DARK)
rc(s,Inches(0.8),Inches(0.9),Inches(0.8),Inches(0.03),ORANGE)
rc(s,Inches(0.8),Inches(1.3),Inches(6.0),Inches(4.5),LGRAY)
tx(s,Inches(1.0),Inches(1.4),Inches(5.6),Inches(0.5),"\u2705 核心模块已跑通",TF,20,True,GREEN)
dones=["Spec \u2192 Code \u2192 Verify \u2192 Diagnose \u2192 Repair 端到端已跑通","规约双重校验（正则 + Dafny resolve）","6 种验证错误的结构化分类与解析","重复错误自动检测 + 强制换策略","HumanEval 前 5 题已完成基线验证"]
for i,item in enumerate(dones):
    tx(s,Inches(1.2),Inches(2.0+i*0.6),Inches(5.3),Inches(0.5),f"\u2713  {item}",BF,15,False,BLACK)
rc(s,Inches(7.0),Inches(1.3),Inches(5.8),Inches(4.5),RGBColor(0xFF,0xF8,0xE1))
tx(s,Inches(7.2),Inches(1.4),Inches(5.4),Inches(0.5),"\u23f3 接下来做",TF,20,True,ORANGE)
todos=["Coordinator Agent \u2014\u2014 退化检测 + 回退","Memory Agent \u2014\u2014 经验向量检索","HumanEval 全量 164 题评测","MBPP 泛化性验证","5 组消融实验"]
for i,item in enumerate(todos):
    tx(s,Inches(7.4),Inches(2.0+i*0.6),Inches(5.1),Inches(0.5),f"\u23f3  {item}",BF,15,False,ORANGE)
rd(s,Inches(0.8),Inches(6.1),Inches(12.0),Inches(0.7),DARK)
tx(s,Inches(1.0),Inches(6.15),Inches(11.6),Inches(0.55),
   "基线数据：HumanEval 前 5 题  端到端通过率 40%（2/5）  平均修复轮次 2.2  简单问题 1 轮收敛 \u2705",BF,15,True,WHITE,PP_ALIGN.CENTER)
print("S11 OK")

# ====== S12 创新点 ======
s=blank(); fs(s,DARK)
rc(s,0,Inches(0.15),W,Inches(0.05),ORANGE)
tx(s,Inches(1),Inches(0.5),Inches(11),Inches(0.6),"三个创新点",TF,34,True,WHITE,PP_ALIGN.CENTER)
rc(s,Inches(5.5),Inches(1.1),Inches(2.3),Inches(0.04),ORANGE)
innvs=[("\u2460 规约-代码协同生成","规约不是事后的补丁，而是生成阶段的约束。规约自带 Dafny 双重校验。",DARK),
       ("\u2461 验证反馈结构化翻译","Dafny 的错误信息 \u2192 分类/定位/修复引导。LLM 不需要自己猜验证器在说什么。",MID),
       ("\u2462 退化感知的修复协议","Coordinator 监控修复质量。越修越差？自动回退 + 换策略。不是固定轮次死循环。",ORANGE)]
for i,(t,d,cl) in enumerate(innvs):
    y=Inches(1.6+i*1.8)
    rc(s,Inches(1),y,Inches(0.9),Inches(1.4),cl)
    rc(s,Inches(2.0),y,Inches(10.7),Inches(1.4),RGBColor(0x2A,0x4A,0x6E))
    tx(s,Inches(2.3),y+Inches(0.15),Inches(10.1),Inches(0.4),t,TF,20,True,WHITE)
    tx(s,Inches(2.3),y+Inches(0.6),Inches(10.1),Inches(0.7),d,BF,16,False,RGBColor(0xBB,0xCC,0xDD),ls=1.4)
rc(s,Inches(1),Inches(6.0),Inches(11.3),Inches(0.8),RGBColor(0x15,0x2D,0x4A))
tx(s,Inches(1.2),Inches(6.05),Inches(10.9),Inches(0.7),
   "与现有工作的最大差异：规约引导 + 结构化反馈翻译，而不是把原始错误文本丢给 LLM 让它自己猜",TF,16,False,RGBColor(0x8A,0xB4,0xF8),PP_ALIGN.CENTER)
print("S12 OK")

# ====== S13 时间线 ======
s=blank()
rc(s,0,0,W,Inches(0.06),DARK)
tx(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.6),"时间线",TF,30,True,DARK)
rc(s,Inches(0.8),Inches(0.9),Inches(0.8),Inches(0.03),ORANGE)
phs=[("第 1-2 月","基础",[("LangGraph 框架","\u2705"),("Dafny 集成","\u2705"),("数据转换","\u2705"),("端到端跑通","\u2705")],DARK),
     ("第 3-4 月","实验",[("Spec/Code Agent","\u2705"),("Diagnose/Repair","\u2705"),("HumanEval 164题","\u23f3"),("消融实验","\u23f3")],MID),
     ("第 5-6 月","完善",[("Coordinator","\u23f3"),("Memory Agent","\u23f3"),("MBPP 泛化","\U0001f4dd"),("论文初稿","\U0001f4dd")],BLUE),
     ("第 7-8 月","投稿",[("论文修改","\U0001f4dd"),("补充实验","\U0001f4dd"),("开源","\U0001f4dd"),("CCF-B投递","\U0001f4dd")],ORANGE)]
for i,(p,t,tasks,cl) in enumerate(phs):
    x=Inches(0.5+i*3.2)
    rc(s,x,Inches(1.3),Inches(3.0),Inches(0.55),cl)
    tx(s,x,Inches(1.32),Inches(3.0),Inches(0.25),p,EF,14,True,WHITE,PP_ALIGN.CENTER)
    tx(s,x,Inches(1.52),Inches(3.0),Inches(0.3),t,TF,14,True,WHITE,PP_ALIGN.CENTER)
    rc(s,x,Inches(1.85),Inches(3.0),Inches(3.5),LGRAY)
    for j,(task,status) in enumerate(tasks):
        tx(s,x+Inches(0.15),Inches(2.0+j*0.6),Inches(2.7),Inches(0.5),f"{status}  {task}",BF,15,False,BLACK if status=="\u2705" else ORANGE)
print("S13 OK")

# ====== S14 谢谢 ======
s=blank(); fs(s,DARK)
rc(s,0,Inches(0.15),W,Inches(0.05),ORANGE)
tx(s,Inches(1),Inches(2.0),Inches(11),Inches(1.0),"谢谢  请老师批评指正",TF,42,True,WHITE,PP_ALIGN.CENTER)
rc(s,Inches(5.5),Inches(3.2),Inches(2.3),Inches(0.04),ORANGE)
bx=s.shapes.add_textbox(Inches(3),Inches(3.6),Inches(7),Inches(1.5)); bx.text_frame.word_wrap=True
for i,(t,fn,fs2,b,c,al) in enumerate([("Tez",TF,22,True,WHITE,PP_ALIGN.CENTER),("软件工程  \u00b7  华东师范大学",BF,16,False,RGBColor(0x8A,0xB4,0xF8),PP_ALIGN.CENTER)]):
    p=bx.text_frame.paragraphs[0] if i==0 else bx.text_frame.add_paragraph()
    p.text=t; p.font.name=fn; p.font.size=Pt(fs2); p.font.bold=b; p.font.color.rgb=c; p.alignment=al
rc(s,0,H-Inches(0.15),W,Inches(0.05),ORANGE)
print("S14 OK")

# ====== Save ======
out = "D:\\codegen-verify\\开题答辩PPT.pptx"
try:
    prs.save(out)
    print(f"\n\u2705 完成！共 {len(prs.slides)} 页 \u2192 {out}")
except:
    out2 = "D:\\codegen-verify\\开题答辩PPT_v4.pptx"
    prs.save(out2)
    print(f"\n\u2705 完成！共 {len(prs.slides)} 页 \u2192 {out2}")

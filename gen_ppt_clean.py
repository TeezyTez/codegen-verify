# -*- coding: utf-8 -*-
"""开题答辩PPT — 简洁版。白底，干净，聚焦内容。"""
import sys; sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation; from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor; from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height

# ====== 颜色方案（简洁 = 少即是多）======
NAVY = RGBColor(0x1B,0x2A,0x4A)      # 深蓝 — 标题
BLUE = RGBColor(0x2B,0x5C,0x9A)       # 中蓝 — 强调
LBLUE = RGBColor(0xE8,0xEF,0xF5)      # 浅蓝 — 背景块
ORANGE = RGBColor(0xD4,0x7B,0x2A)     # 橙色 — 高亮
GRAY = RGBColor(0x66,0x66,0x66)        # 灰色 — 正文
LGRAY = RGBColor(0xF5,0xF5,0xF5)      # 浅灰 — 轻背景
WHITE = RGBColor(0xFF,0xFF,0xFF)
DARK = RGBColor(0x33,0x33,0x33)
GREEN = RGBColor(0x2E,0x7D,0x32)
RED = RGBColor(0xCC,0x33,0x33)
TF = 'Microsoft YaHei'               # 中文字体
EF = 'Calibri'                        # 英文/数字

def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

def add_rect(s, l, t, w, h, fill, alpha=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if alpha is not None: sh.fill.fore_color.brightness = alpha
    sh.line.fill.background()
    return sh

def add_round_rect(s, l, t, w, h, fill):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh

def add_border_rect(s, l, t, w, h, fill, border):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border; sh.line.width = Pt(1)
    return sh

def txt(s, l, t, w, h, text, font=TF, size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT, spacing=1.3):
    bx = s.shapes.add_textbox(l, t, w, h)
    bx.text_frame.word_wrap = True
    p = bx.text_frame.paragraphs[0]
    p.text = text
    p.font.name = font; p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = color; p.alignment = align
    p.space_after = Pt(0); p.space_before = Pt(0)
    from pptx.oxml.ns import qn
    pPr = p._p.get_or_add_pPr()
    ln = pPr.makeelement(qn('a:lnSpc'), {})
    sp = ln.makeelement(qn('a:spcPct'), {'val': str(int(spacing * 100000))})
    ln.append(sp); pPr.append(ln)
    return bx

def multi_txt(s, l, t, w, h, lines, spacing=1.3):
    """lines = [(text, font, size, bold, color, align), ...]"""
    bx = s.shapes.add_textbox(l, t, w, h)
    bx.text_frame.word_wrap = True
    from pptx.oxml.ns import qn
    for i, (text, fn, sz, b, c, al) in enumerate(lines):
        p = bx.text_frame.paragraphs[0] if i == 0 else bx.text_frame.add_paragraph()
        p.text = text; p.font.name = fn; p.font.size = Pt(sz)
        p.font.bold = b; p.font.color.rgb = c; p.alignment = al
        p.space_after = Pt(0); p.space_before = Pt(0)
        pPr = p._p.get_or_add_pPr()
        ln = pPr.makeelement(qn('a:lnSpc'), {})
        sp = ln.makeelement(qn('a:spcPct'), {'val': str(int(spacing * 100000))})
        ln.append(sp); pPr.append(ln)
    return bx

def slide_header(s, title, subtitle=None):
    """统一页面顶部：标题栏 + 分隔线"""
    add_rect(s, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.04), NAVY)
    txt(s, Inches(0.8), Inches(0.55), Inches(10), Inches(0.5), title, TF, 24, True, NAVY)
    if subtitle:
        txt(s, Inches(0.8), Inches(1.0), Inches(10), Inches(0.35), subtitle, TF, 13, False, GRAY)

# ============================================================
# S1: 封面
# ============================================================
s = blank()
add_rect(s, 0, 0, W, H, NAVY)  # 全页深蓝
txt(s, Inches(1), Inches(1.0), Inches(11.3), Inches(0.4),
    "华东师范大学  ·  软件工程  ·  硕士论文开题答辩", TF, 16, False, RGBColor(0xAA,0xBB,0xDD), PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(2.2), Inches(11.3), Inches(1.2),
    "基于形式化规约的 LLM\n代码生成与自修复方法研究", TF, 40, True, WHITE, PP_ALIGN.CENTER, 1.3)
add_rect(s, Inches(5.5), Inches(3.8), Inches(2.3), Inches(0.04), ORANGE)
multi_txt(s, Inches(3.5), Inches(4.1), Inches(6.3), Inches(2.0), [
    ("Tez", TF, 22, True, WHITE, PP_ALIGN.CENTER),
    ("导师：XXX 教授", TF, 16, False, RGBColor(0xBB,0xCC,0xDD), PP_ALIGN.CENTER),
    ("2026年6月", TF, 14, False, RGBColor(0x88,0xAA,0xCC), PP_ALIGN.CENTER),
], spacing=1.6)
add_rect(s, 0, Inches(7.35), W, Inches(0.15), ORANGE)
print("S1 OK")

# ============================================================
# S2: 目录
# ============================================================
s = blank()
slide_header(s, "目录")
items = [("01", "研究背景"), ("02", "方案与架构"), ("03", "初步实验"),
         ("04", "研究进度"), ("05", "创新点")]
for i, (num, label) in enumerate(items):
    x = Inches(1.5) + i * Inches(2.2)
    y = Inches(2.0)
    # 编号圆
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.7), Inches(0.7))
    c.fill.solid(); c.fill.fore_color.rgb = NAVY; c.line.fill.background()
    tf = c.text_frame; tf.paragraphs[0].text = num
    tf.paragraphs[0].font.name = EF; tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    txt(s, x + Inches(0.85), y + Inches(0.15), Inches(1.5), Inches(0.4), label, TF, 16, True, NAVY)
# 底部提示
txt(s, Inches(1), Inches(5.0), Inches(11.3), Inches(0.4),
    "核心思路：用形式化规约（Dafny）作为 LLM 代码生成的独立裁判，用结构化反馈引导自修复",
    TF, 13, False, GRAY, PP_ALIGN.CENTER)
print("S2 OK")

# ============================================================
# S3: 研究背景
# ============================================================
s = blank()
slide_header(s, "研究背景", "LLM 写代码很强，但正确性怎么保证？")

# 左侧：现状
add_round_rect(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.2), LBLUE)
txt(s, Inches(1.1), Inches(1.6), Inches(5.0), Inches(0.35), "现状：生成能力已很强", TF, 16, True, NAVY)
txt(s, Inches(1.1), Inches(2.0), Inches(5.0), Inches(1.5),
    "Codex (2021) → ReflexiCoder (2026)：HumanEval 从 28.8% 涨到 94.5%\n"
    "GitHub Copilot 被百万组织使用\n\n"
    "但——测试覆盖是假象：每题仅 7.8 个测试用例\n"
    "扩到 774 个后正确率从 85% 掉到 33% (Yu 2025)",
    TF, 14, False, DARK, spacing=1.5)

# 右侧：问题
add_round_rect(s, Inches(6.8), Inches(1.5), Inches(5.7), Inches(2.2), RGBColor(0xFD,0xF0,0xF0))
txt(s, Inches(7.1), Inches(1.6), Inches(5.2), Inches(0.35), "问题：自修复有系统缺陷", TF, 16, True, RED)
txt(s, Inches(7.1), Inches(2.0), Inches(5.2), Inches(1.5),
    "自审漏检率 31.7% (Reddy 2026)\n"
    "误拒率飙到 73%（越详细提示越惨）(Jin 2026)\n"
    "断言错误修复率仅 45% (Arimbur 2026)\n\n"
    "LLM 同时当考生和阅卷老师 → 认知盲区无法克服",
    TF, 14, False, DARK, spacing=1.5)

# 底部：解法方向
add_round_rect(s, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.5), LGRAY)
txt(s, Inches(1.1), Inches(4.3), Inches(11.2), Inches(0.35), "解法方向：引入独立的形式化验证器", TF, 16, True, NAVY)
txt(s, Inches(1.1), Inches(4.75), Inches(11.2), Inches(1.8),
    "规约（requires / ensures）是对「什么算对」的精确回答 → 验证器给出数学级别的正确性保证\n"
    "但手写规约门槛高，验证器输出对 LLM 不友好\n\n"
    "关键思路：\n"
    "  LLM 写规约 → LLM 写代码 → Dafny 当裁判 → Diagnose Agent 翻译裁判的话 → LLM 自己修\n\n"
    "现有问题：规约生成与代码生成分离、验证反馈没结构、修复缺少退化控制",
    TF, 14, False, DARK, spacing=1.4)
print("S3 OK")

# ============================================================
# S4: 研究目标
# ============================================================
s = blank()
slide_header(s, "研究目标", "四件事：规约 → 生成 → 翻译 → 修复")

goals = [
    ("内容一", "规约自动生成", "NL → Dafny requires/ensures\n自带语法合法性校验", NAVY),
    ("内容二", "规约感知代码生成", "规约作为第一优先级约束\n嵌入代码生成 Prompt", BLUE),
    ("内容三", "验证反馈结构化翻译", "Dafny 原始错误 → 分类/定位/修复引导\n让 LLM 直接使用", RGBColor(0x3A,0x7C,0xBD)),
    ("内容四", "退化感知迭代修复", "Coordinator 监控修复质量\n越修越差 → 自动回退 + 换策略", ORANGE),
]
for i, (n, t, d, cl) in enumerate(goals):
    x = Inches(0.5) + i * Inches(3.2)
    add_round_rect(s, x, Inches(1.6), Inches(3.0), Inches(2.6), LBLUE) if i < 2 else add_round_rect(s, x, Inches(1.6), Inches(3.0), Inches(2.6), RGBColor(0xF5,0xEE,0xE8) if i==3 else RGBColor(0xE8,0xF0,0xF8))
    add_rect(s, x + Inches(0.1), Inches(1.7), Inches(2.8), Inches(0.04), cl)
    txt(s, x + Inches(0.15), Inches(1.85), Inches(2.7), Inches(0.3), n, TF, 11, True, cl, PP_ALIGN.CENTER)
    txt(s, x + Inches(0.15), Inches(2.15), Inches(2.7), Inches(0.3), t, TF, 15, True, NAVY, PP_ALIGN.CENTER)
    txt(s, x + Inches(0.15), Inches(2.6), Inches(2.7), Inches(1.3), d, TF, 14, False, DARK, PP_ALIGN.CENTER, 1.5)

# 底部架构总览
add_round_rect(s, Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.5), LGRAY)
txt(s, Inches(1.0), Inches(4.9), Inches(11.3), Inches(0.35), "六 Agent 架构总览", TF, 15, True, NAVY)
txt(s, Inches(1.0), Inches(5.3), Inches(11.3), Inches(0.8),
    "问题描述 → [Spec Agent] → [Code Agent] → [Dafny Verify] → [Diagnose Agent] → [Repair Agent] → 循环\n"
    "                               [Coordinator Agent] 路由决策  ·  [Memory Agent] 检索经验（规划中）",
    TF, 14, False, DARK, PP_ALIGN.CENTER, 1.5)
print("S4 OK")

# ============================================================
# S5: 技术方案 — 反馈翻译
# ============================================================
s = blank()
slide_header(s, "核心技术：验证反馈的结构化翻译", "Dafny 说的是什么 → Diagnose Agent 翻成 LLM 听得懂的")

# Dafny 原始输出
add_round_rect(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.8), RGBColor(0xFD,0xED,0xED))
txt(s, Inches(1.0), Inches(1.6), Inches(5.1), Inches(0.35), "Dafny 原始输出", TF, 16, True, RED)
txt(s, Inches(1.0), Inches(2.05), Inches(5.1), Inches(2.0),
    "stdin.dfy(10,4): Error: A postcondition\n"
    "might not hold on this return path.\n"
    "stdin.dfy(4,10): Related location:\n"
    "This is the postcondition that might not hold.",
    EF, 14, False, DARK, spacing=1.5)

# 箭头
ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.5), Inches(2.4), Inches(0.6), Inches(0.35))
ar.fill.solid(); ar.fill.fore_color.rgb = ORANGE; ar.line.fill.background()

# Diagnose 输出
add_round_rect(s, Inches(7.3), Inches(1.5), Inches(5.2), Inches(2.8), RGBColor(0xE8,0xF5,0xE9))
txt(s, Inches(7.5), Inches(1.6), Inches(4.8), Inches(0.35), "Diagnose Agent 翻译后", TF, 16, True, GREEN)
txt(s, Inches(7.5), Inches(2.05), Inches(4.8), Inches(2.0),
    "错误类型：postcondition_violation\n"
    "位置：第10行第4列\n"
    "关联规约：ensures result >= x\n"
    "修复建议：检查该分支上 result 的计算",
    TF, 14, False, DARK, spacing=1.5)

# 底部：三种错误策略
add_round_rect(s, Inches(0.8), Inches(4.7), Inches(11.7), Inches(2.0), LGRAY)
txt(s, Inches(1.0), Inches(4.8), Inches(11.3), Inches(0.35), "错误类型 → 修复策略映射", TF, 15, True, NAVY)
errs = [("postcondition", "修正函数逻辑或条件", NAVY),
        ("invariant", "调整循环不变量或循环体", BLUE),
        ("syntax / type", "直接修正语法或类型标注", ORANGE)]
for i, (et, fix, cl) in enumerate(errs):
    x = Inches(1.0) + i * Inches(3.8)
    add_rect(s, x, Inches(5.3), Inches(3.4), Inches(0.45), cl)
    txt(s, x, Inches(5.32), Inches(3.4), Inches(0.4), et, TF, 14, True, WHITE, PP_ALIGN.CENTER)
    txt(s, x, Inches(5.85), Inches(3.4), Inches(0.4), fix, TF, 13, False, DARK, PP_ALIGN.CENTER)

# 附：assert 桥接
add_rect(s, Inches(1.0), Inches(6.35), Inches(11.5), Inches(0.5), LBLUE)
txt(s, Inches(1.2), Inches(6.38), Inches(11.1), Inches(0.45),
    "助攻：_inject_nested_loop_assert() 自动检测嵌套循环，在 i:=i+1 前注入 assert 桥接，解决 Dafny 量词实例化失败的问题",
    TF, 12, False, NAVY)
print("S5 OK")

# ============================================================
# S6: 实验方案
# ============================================================
s = blank()
slide_header(s, "实验方案", "怎么验证我做得对不对")

# 数据集
add_round_rect(s, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.7), LBLUE)
txt(s, Inches(1.0), Inches(1.35), Inches(11.3), Inches(0.55),
    "数据集：HumanEval（164 题，主实验）  |  MBPP（974 题，泛化性验证）  |  Proving the Coding Interview（27 题，参照）",
    TF, 14, False, NAVY, PP_ALIGN.CENTER)

# 基线
add_round_rect(s, Inches(0.8), Inches(2.3), Inches(5.7), Inches(1.8), LGRAY)
txt(s, Inches(1.0), Inches(2.4), Inches(5.3), Inches(0.35), "基线对比", TF, 15, True, NAVY)
bls = "Direct Generation（绝对底线）\nSelf-Debug（执行信号修复）\nReflexion（经验驱动修复）\nReflexiCoder（最先进模型）\n单 Agent Pipeline（消融基线）"
txt(s, Inches(1.0), Inches(2.85), Inches(5.3), Inches(1.2), bls, TF, 13, False, DARK, spacing=1.5)

# 指标
add_round_rect(s, Inches(6.8), Inches(2.3), Inches(5.7), Inches(1.8), LGRAY)
txt(s, Inches(7.0), Inches(2.4), Inches(5.3), Inches(0.35), "评估指标", TF, 15, True, NAVY)
idx = "Verif@1（首轮通过率）\nVerif@k（k 轮通过率）\nSpecCompile（规约编译通过率）\nAvgRounds（平均轮次）\nDegradationRate（退化率）"
txt(s, Inches(7.0), Inches(2.85), Inches(5.3), Inches(1.2), idx, TF, 13, False, DARK, spacing=1.5)

# 消融实验
add_round_rect(s, Inches(0.8), Inches(4.4), Inches(11.7), Inches(2.5), LGRAY)
txt(s, Inches(1.0), Inches(4.5), Inches(11.3), Inches(0.35), "六组消融实验", TF, 15, True, NAVY)
abs_ = [
    ("A1 规约有无", "规约对首轮通过率的贡献"),
    ("A2 反馈结构", "结构化翻译 vs Dafny 原始输出"),
    ("A3 Agent 分离", "单 Agent vs 多 Agent"),
    ("A4 轮次上限", "1/3/5 轮效果对比"),
    ("A5 记忆有无", "Memory Agent 的贡献"),
    ("A6 模型对比", "DeepSeek / GPT-4o / Qwen3"),
]
for i, (n, d) in enumerate(abs_):
    x = Inches(0.8) + (i % 3) * Inches(3.9)
    y = Inches(5.0) + (i // 3) * Inches(0.9)
    add_round_rect(s, x, y, Inches(3.6), Inches(0.75), WHITE)
    txt(s, x + Inches(0.1), y + Inches(0.05), Inches(3.4), Inches(0.3), n, TF, 14, True, ORANGE)
    txt(s, x + Inches(0.1), y + Inches(0.38), Inches(3.4), Inches(0.3), d, TF, 12, False, GRAY)

print("S6 OK")

# ============================================================
# S7: 初步实验成果
# ============================================================
s = blank()
slide_header(s, "初步实验成果", "HumanEval 前 5 题基线（DeepSeek 模型，max_rounds=3）")

# 成果表
table_data = [
    ("题目", "函数", "通过?", "轮次", "耗时"),
    ("HumanEval/0", "has_close_elements", "✅ 是", "1", "12s"),
    ("HumanEval/1", "separate_paren_groups", "❌ 否", "3", "66s"),
    ("HumanEval/2", "truncate_number", "✅ 是", "1", "7s"),
    ("HumanEval/3", "below_zero", "❌ 否", "3", "45s"),
    ("HumanEval/4", "mean_abs_deviation", "❌ 否", "3", "47s"),
]

# 手动绘制表格
col_w = [Inches(2.0), Inches(3.5), Inches(1.5), Inches(1.2), Inches(1.2)]
col_xs = [Inches(0.8)]
for w in col_w[:-1]:
    col_xs.append(col_xs[-1] + w)
total_w = sum(col_w)

# 表头
for j, (head, cx) in enumerate(zip(table_data[0], col_xs)):
    add_rect(s, cx, Inches(1.4), col_w[j], Inches(0.45), NAVY)
    txt(s, cx + Inches(0.1), Inches(1.42), col_w[j] - Inches(0.2), Inches(0.4), head, TF, 14, True, WHITE, PP_ALIGN.CENTER)

for i, row in enumerate(table_data[1:]):
    y = Inches(1.85) + i * Inches(0.5)
    bg = WHITE if i % 2 == 0 else LGRAY
    for j, (cell, cx) in enumerate(zip(row, col_xs)):
        add_border_rect(s, cx, y, col_w[j], Inches(0.5), bg, RGBColor(0xDD,0xDD,0xDD))
        cl = GREEN if cell == "✅ 是" else (RED if cell == "❌ 否" else DARK)
        txt(s, cx + Inches(0.1), y + Inches(0.05), col_w[j] - Inches(0.2), Inches(0.4),
            cell, TF, 14, False, cl, PP_ALIGN.CENTER)

# 汇总
add_round_rect(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(1.0), LBLUE)
txt(s, Inches(1.0), Inches(4.65), Inches(11.3), Inches(0.8),
    "端到端通过率：40%（2/5）  |  成功案例均一轮通过\n"
    "平均修复轮次：2.2  |  规约语法编译通过率：100%\n"
    "瓶颈：未通过的 3 题问题在循环不变量（invariant）的精确性上",
    TF, 14, False, NAVY, spacing=1.5)

# 进展总结
add_round_rect(s, Inches(0.8), Inches(5.9), Inches(5.7), Inches(1.0), RGBColor(0xE8,0xF5,0xE9))
txt(s, Inches(1.0), Inches(5.95), Inches(5.3), Inches(0.8),
    "✅ 已完成：核心 5 Agent 跑通 · 规约双重校验\n"
    "· 结构化错误解析 · 重复错误换策略 · 65+ 篇论文精读",
    TF, 13, False, GREEN, spacing=1.4)

add_round_rect(s, Inches(6.8), Inches(5.9), Inches(5.7), Inches(1.0), RGBColor(0xFF,0xF8,0xE1))
txt(s, Inches(7.0), Inches(5.95), Inches(5.3), Inches(0.8),
    "⏳ 待完成：HumanEval 全 164 题 · NLR2VC-60 评测\n"
    "· Coordinator 退化检测 · Memory Agent · 消融实验",
    TF, 13, False, ORANGE, spacing=1.4)
print("S7 OK")

# ============================================================
# S8: 时间线
# ============================================================
s = blank()
slide_header(s, "研究计划", "8 个月 → 论文 + 系统")

phases = [
    ("第 1-2 月", "目前已完成", NAVY, ["LangGraph 框架搭建 ✅", "Dafny 集成 + 结构化解析 ✅", "数据脚本 ✅", "端到端 5 题跑通 ✅"]),
    ("第 3-4 月", "全面实验", BLUE, ["HumanEval 全 164 题", "NL2VC-60 基准评测", "Coordinator / Memory Agent", "对比实验 + 消融实验"]),
    ("第 5-6 月", "论文撰写", RGBColor(0x3A,0x7C,0xBD), ["MBPP 泛化性验证", "多模型对比", "论文第一稿", "补充 Reviewer 实验"]),
    ("第 7-8 月", "修改投稿", ORANGE, ["论文修改完善", "开源代码整理发布", "CCF-B 类会议/期刊投稿", "学位论文初稿"]),
]

for i, (phase, desc, cl, tasks) in enumerate(phases):
    x = Inches(0.5) + i * Inches(3.2)
    add_rect(s, x, Inches(1.3), Inches(3.0), Inches(0.55), cl)
    txt(s, x, Inches(1.32), Inches(3.0), Inches(0.25), phase, EF, 14, True, WHITE, PP_ALIGN.CENTER)
    txt(s, x, Inches(1.52), Inches(3.0), Inches(0.25), desc, TF, 12, True, WHITE, PP_ALIGN.CENTER)
    add_round_rect(s, x, Inches(1.85), Inches(3.0), Inches(4.0), LGRAY)
    for j, task in enumerate(tasks):
        icon = "✅" if "✅" in task else "⏳"
        clean = task.replace(" ✅", "").replace(" ⏳", "")
        txt(s, x + Inches(0.15), Inches(2.0 + j * 0.7), Inches(2.7), Inches(0.5),
            f"{icon}  {clean}", TF, 14, False, GREEN if "✅" in task else DARK)

# 底部预期成果
add_round_rect(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.6), NAVY)
txt(s, Inches(1.0), Inches(6.35), Inches(11.3), Inches(0.5),
    "预期成果：原型系统 | HumanEval Dafny 规约增强集 | CCF-B 类会议/期刊论文 1-2 篇 | 开源代码",
    TF, 14, False, WHITE, PP_ALIGN.CENTER)
print("S8 OK")

# ============================================================
# S9: 创新点
# ============================================================
s = blank()
slide_header(s, "三个创新点")

innovs = [
    ("① 规约-代码协同生成",
     "规约不是事后的补丁，而是生成阶段的约束。Spec Agent 产出的规约经过 Dafny 合法性校验后，直接约束 Code Agent 的生成过程。多候选 + 规约过滤 → 生成阶段就内置了验证意识。"),
    ("② 验证反馈的结构化翻译",
     "Dafny 的原始错误信息 → Diagnose Agent 翻译为分类+定位+修复引导。LLM 不需要自己猜验证器在说什么。这是对 Olausson (2023)「反馈质量瓶颈」的直接回应。"),
    ("③ 退化感知的迭代修复协议",
     "Coordinator 监控修复质量走势。检测到退化（score 下降 20%）→ 自动回退到最佳版本 + 切换保守策略。不是固定 3 轮的死循环，而是有质量意识的自适应控制。"),
]

for i, (title, desc) in enumerate(innovs):
    y = Inches(1.3 + i * 1.8)
    add_rect(s, Inches(0.8), y, Inches(0.08), Inches(1.5), [NAVY, BLUE, ORANGE][i])
    txt(s, Inches(1.1), y, Inches(11.0), Inches(0.35), title, TF, 18, True, NAVY)
    txt(s, Inches(1.1), y + Inches(0.4), Inches(11.0), Inches(1.0), desc, TF, 14, False, DARK, spacing=1.5)

# 一句话定位
add_round_rect(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.8), NAVY)
txt(s, Inches(1.0), Inches(6.1), Inches(11.3), Inches(0.6),
    "一句话定位：规约引导 + 结构化反馈翻译，不是把原始错误文本丢给 LLM 让它自己猜",
    TF, 15, False, WHITE, PP_ALIGN.CENTER, 1.3)
print("S9 OK")

# ============================================================
# S10: 谢谢
# ============================================================
s = blank()
add_rect(s, 0, 0, W, H, NAVY)
txt(s, Inches(1), Inches(2.5), Inches(11.3), Inches(0.8), "谢谢  请老师批评指正", TF, 40, True, WHITE, PP_ALIGN.CENTER)
add_rect(s, Inches(5.5), Inches(3.5), Inches(2.3), Inches(0.04), ORANGE)
txt(s, Inches(1), Inches(3.8), Inches(11.3), Inches(0.4), "Tez  ·  软件工程  ·  华东师范大学", TF, 18, False, RGBColor(0xAA,0xBB,0xDD), PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(4.3), Inches(11.3), Inches(0.4), "2026年6月", TF, 14, False, RGBColor(0x88,0xAA,0xCC), PP_ALIGN.CENTER)
add_rect(s, 0, Inches(7.35), W, Inches(0.15), ORANGE)
print("S10 OK")

# ====== Save ======
out = "D:\\codegen-verify\\开题答辩PPT_简洁版.pptx"
try:
    prs.save(out)
    print(f"\n✅ 完成！共 {len(prs.slides)} 页 → {out}")
except Exception as e:
    print(f"\n❌ 保存失败: {e}")

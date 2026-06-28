"""
生成开题答辩PPT
基于 python-pptx
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ========== 颜色方案 ==========
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)     # 深蓝（标题背景）
MID_BLUE = RGBColor(0x2C, 0x5F, 0x8A)      # 中蓝
LIGHT_BLUE = RGBColor(0x3A, 0x7C, 0xBD)    # 亮蓝
ACCENT_ORANGE = RGBColor(0xE8, 0x6C, 0x00)  # 强调橙
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
RED = RGBColor(0xCC, 0x33, 0x33)
GREEN = RGBColor(0x2E, 0x7D, 0x32)

# ========== 字体 ==========
TITLE_FONT = 'SimHei'   # 黑体
BODY_FONT = 'SimSun'    # 宋体
EN_FONT = 'Calibri'

# ========== 创建演示文稿 ==========
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ========== 工具函数 ==========

def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_bg_rect(slide, color=DARK_BLUE):
    """添加全屏背景色矩形"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, W, H
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_rect(slide, left, top, width, height, color, alpha=None):
    """添加矩形"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        solidFill = shape.fill._fill
        srgb = solidFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgb is not None:
            alpha_elem = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha_elem.set('val', str(int(alpha * 1000)))
    return shape


def add_textbox(slide, left, top, width, height, text, font_name=BODY_FONT,
                font_size=18, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                line_spacing=1.3, anchor=MSO_ANCHOR.TOP, italic=False):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None
    
    # 设置垂直对齐
    tf = txBox.text_frame
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    
    # line spacing
    from pptx.oxml.ns import qn
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
    spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
    lnSpc.append(spcPct)
    pPr.append(lnSpc)
    
    # 垂直对齐
    txBox.text_frame.paragraphs[0].alignment = alignment
    try:
        txBox.text_frame.paragraphs[0].font._element.attrib
    except:
        pass
    
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines_with_format,
                          anchor=MSO_ANCHOR.TOP, line_spacing=1.3):
    """
    多行文本框
    lines_with_format: list of (text, font_name, font_size, bold, color, alignment)
    每行独立格式
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    tf = txBox.text_frame
    
    from pptx.oxml.ns import qn
    # 设置 line spacing
    pPr = tf.paragraphs[0]._pPr
    if pPr is None:
        pPr = tf.paragraphs[0]._p.get_or_add_pPr()
    lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
    spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
    lnSpc.append(spcPct)
    pPr.append(lnSpc)
    
    for i, (text, fn, fs, bld, clr, align) in enumerate(lines_with_format):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = text
        p.font.name = fn
        p.font.size = Pt(fs)
        p.font.bold = bld
        p.font.color.rgb = clr
        p.alignment = align
        p.space_before = Pt(2)
        p.space_after = Pt(2)
    
    return txBox


def add_bullet_textbox(slide, left, top, width, height, items, font_name=BODY_FONT,
                       font_size=16, color=BLACK, bullet_char="●", line_spacing=1.3):
    """带项目符号的文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    tf = txBox.text_frame
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        full_text = f"{bullet_char} {item}"
        p.text = full_text
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        
        from pptx.oxml.ns import qn
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
        spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
        lnSpc.append(spcPct)
        pPr.append(lnSpc)
    
    return txBox


def add_arrow_shape(slide, left, top, width, height, color=LIGHT_BLUE):
    """添加右箭头"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_chevron(slide, left, top, width, height, color=MID_BLUE):
    """添加 V形"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_flow_box(slide, left, top, width, height, text, bg_color=MID_BLUE, 
                 text_color=WHITE, font_size=12, font_name=TITLE_FONT):
    """添加流程框"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.name = font_name
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.color.rgb = text_color
    tf.paragraphs[0].font.bold = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_page_title(slide, title_text, subtitle_text=None):
    """在顶部添加统一的页面标题"""
    # 顶部蓝条
    add_rect(slide, 0, 0, W, Inches(0.08), DARK_BLUE)
    
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                title_text, TITLE_FONT, 28, True, DARK_BLUE, PP_ALIGN.LEFT)
    
    if subtitle_text:
        add_textbox(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.4),
                    subtitle_text, BODY_FONT, 14, False, GRAY, PP_ALIGN.LEFT)
    
    # 分隔线
    add_rect(slide, Inches(0.8), Inches(1.15), Inches(1.5), Inches(0.04), ACCENT_ORANGE)


# ================================================================
# Slide 1: 封面
# ================================================================
slide = add_blank_slide()
add_bg_rect(slide, DARK_BLUE)

# 顶部装饰线
add_rect(slide, 0, Inches(0.15), W, Inches(0.06), ACCENT_ORANGE)

# 学校名称
add_textbox(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.5),
            "华东师范大学", TITLE_FONT, 20, False, RGBColor(0x8A, 0xB4, 0xF8),
            PP_ALIGN.CENTER)

# 主标题
add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
            "基于形式化契约引导的 LLM 代码生成\n与可验证自修复方法研究",
            TITLE_FONT, 36, True, WHITE, PP_ALIGN.CENTER, line_spacing=1.4)

# 分隔线
add_rect(slide, Inches(5.5), Inches(3.3), Inches(2.3), Inches(0.04), ACCENT_ORANGE)

# 信息
info_lines = [
    ("硕士论文开题答辩", BODY_FONT, 18, False, RGBColor(0xCC, 0xD5, 0xE8), PP_ALIGN.CENTER),
    ("", BODY_FONT, 10, False, WHITE, PP_ALIGN.CENTER),
    ("答辩人：Tez", BODY_FONT, 16, False, WHITE, PP_ALIGN.CENTER),
    ("专  业：软件工程", BODY_FONT, 16, False, WHITE, PP_ALIGN.CENTER),
    ("导  师：XXX 教授", BODY_FONT, 16, False, WHITE, PP_ALIGN.CENTER),
]
add_multiline_textbox(slide, Inches(3), Inches(3.6), Inches(7), Inches(2.5), info_lines)

# 底部装饰线
add_rect(slide, 0, H - Inches(0.15), W, Inches(0.06), ACCENT_ORANGE)

# 日期
add_textbox(slide, Inches(1), H - Inches(0.6), Inches(11), Inches(0.4),
            "2026年5月", BODY_FONT, 14, False, RGBColor(0x8A, 0xB4, 0xF8),
            PP_ALIGN.CENTER)

print("封面 ✓")


# ================================================================
# Slide 2: 目录
# ================================================================
slide = add_blank_slide()
add_page_title(slide, "目  录")

items = [
    ("01", "研究背景与问题", "为什么做这个研究？"),
    ("02", "相关工作与关键洞察", "现有方法的边界在哪？"),
    ("03", "研究目标与内容", "我要做什么？"),
    ("04", "技术路线与系统架构", "怎么做？"),
    ("05", "实验方案与预期成果", "怎么验证？"),
    ("06", "创新点与研究计划", "亮点与时间线"),
]

x_start = Inches(1.2)
y_start = Inches(1.6)
x_gap = Inches(3.5)
y_gap = Inches(1.6)

for i, (num, title, desc) in enumerate(items):
    col = i % 3
    row = i // 3
    x = x_start + col * x_gap
    y = y_start + row * y_gap
    
    # 编号圆圈
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x, y, Inches(0.6), Inches(0.6)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = DARK_BLUE if row == 0 else MID_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.name = EN_FONT
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 标题
    add_textbox(slide, x + Inches(0.8), y - Inches(0.05), Inches(2.5), Inches(0.4),
                title, TITLE_FONT, 16, True, DARK_BLUE)
    # 描述
    add_textbox(slide, x + Inches(0.8), y + Inches(0.35), Inches(2.5), Inches(0.3),
                desc, BODY_FONT, 11, False, GRAY)

print("目录 ✓")


# ================================================================
# Slide 3: 研究背景 — LLM代码生成的进步与正确性鸿沟
# ================================================================
slide = add_blank_slide()
add_page_title(slide, "1.1 研究背景", "LLM代码生成能力的飞跃，但正确性保障严重滞后")

# 左侧：能力进步
add_textbox(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.4),
            "▎LLM代码生成能力的飞速提升", TITLE_FONT, 18, True, DARK_BLUE)

code_gen_items = [
    "Codex (2021): HumanEval Pass@1 = 28.8%",
    "Code Llama 70B (2023): Pass@1 = 67%",
    "ReflexiCoder (2026): Pass@1 = 94.51%",
    "GitHub Copilot: 超100万组织采用，每天生成数十亿行代码",
]
add_bullet_textbox(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(2.5),
                   code_gen_items, BODY_FONT, 14, BLACK)

# 右侧：正确性鸿沟
add_textbox(slide, Inches(7.2), Inches(1.4), Inches(5.5), Inches(0.4),
            "▎但正确性保障存在根本性鸿沟", TITLE_FONT, 18, True, RED)

gap_items = [
    "LLM 能 '生成' 代码，但本质上不 '理解' 执行语义",
    "Codex 论文揭示：模型预测自身输出正确性的准确率仅略高于随机",
    "生成代码可能包含隐藏的逻辑错误、边界条件遗漏、安全漏洞",
    "生成能力接近天花板，关注点转向 '如何保证正确性'",
]
add_bullet_textbox(slide, Inches(7.2), Inches(1.9), Inches(5.5), Inches(2.5),
                   gap_items, BODY_FONT, 14, BLACK)

# 底部放大引用
add_rect(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.8), LIGHT_GRAY)
add_textbox(slide, Inches(1.2), Inches(4.95), Inches(11), Inches(0.4),
            "💡 核心问题", TITLE_FONT, 16, True, RED)
add_textbox(slide, Inches(1.2), Inches(5.4), Inches(11), Inches(1.0),
            "当代码生成能力接近 '天花板'，我们能否保证生成的代码不仅仅是 '看起来对'，而是 '数学意义上正确'？\n"
            "形式化验证提供了答案：如果代码满足形式化规约，其正确性有数学证明。",
            BODY_FONT, 14, False, BLACK, PP_ALIGN.LEFT, 1.4)

print("背景1 ✓")


# ================================================================
# Slide 4: 三个可靠性危机
# ================================================================
slide = add_blank_slide()
add_page_title(slide, "1.2 三个可靠性危机", "现有自修复范式的根本性缺陷")

# 三个危机卡片
crises = [
    ("危机一：自审查双向失效", "🔍",
     [
         "Reddy et al. (2026)：LLM 对自身输出的语义错误漏检率达 31.7%",
         "Jin et al. (2026)：正确代码被误判为 '不符合需求' 的比例达 73%",
         "→ 模型在 '漏放' 与 '误拒' 两个方向上同时失效",
     ]),
    ("危机二：反馈质量瓶颈", "🧪",
     [
         "Olausson et al. (2023)：人类反馈优于GPT-4自反馈1.58倍",
         "Arimbur et al. (2026)：断言错误修复率仅45%（修复上限）",
         "Ruiz et al. (2026)：指令微调反而削弱了修复灵活性",
     ]),
    ("危机三：缺乏形式化保障", "🛡️",
     [
         "测试执行反馈 → 覆盖不全，非通过的路径不可见",
         "LLM内省反馈 → 共享生成者的认知盲区",
         "缺少能提供数学级正确性保证的验证手段",
     ]),
]

for i, (title, icon, items) in enumerate(crises):
    x = Inches(0.8) + i * Inches(4.1)
    y = Inches(1.6)
    w = Inches(3.8)
    
    # 卡片背景
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(5.0)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_GRAY
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card.line.width = Pt(1)
    
    # 标题
    add_textbox(slide, x + Inches(0.3), y + Inches(0.2), w - Inches(0.6), Inches(0.5),
                f"{icon} {title}", TITLE_FONT, 14, True, DARK_BLUE, PP_ALIGN.CENTER)
    
    # 分隔
    add_rect(slide, x + Inches(0.5), y + Inches(0.8), w - Inches(1.0), Inches(0.02), ACCENT_ORANGE)
    
    # 内容
    add_bullet_textbox(slide, x + Inches(0.3), y + Inches(1.0), w - Inches(0.6), Inches(3.8),
                       items, BODY_FONT, 12, BLACK, "▸", 1.3)

# 底部结论
add_rect(slide, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.5), DARK_BLUE)
add_textbox(slide, Inches(1.0), Inches(6.82), Inches(11.3), Inches(0.4),
            "共性问题：自修复依赖的反馈信号本身不可靠 → 需要引入形式化验证提供可证明的正确性保证",
            TITLE_FONT, 14, True, WHITE, PP_ALIGN.CENTER)

print("三个危机 ✓")


# ================================================================
# Slide 5: 形式化方案的核心优势
# ================================================================
slide = add_blank_slide()
add_page_title(slide, "1.3 为什么选择形式化方案？", "形式化规约提供了精确、无歧义的代码正确性标准")

# 对比表格 - 用形状模拟
headers = ["对比维度", "测试执行驱动", "LLM自审查驱动", "形式化验证驱动（本方案）"]
rows = [
    ["正确性标准", "测试用例覆盖", "LLM 内省判断", "数学证明（全覆盖）"],
    ["反馈精度", "\"测试不通过\"（黑盒）", "\"可能有问题\"（模糊）", "结构化错误定位（行/列/类型）"],
    ["反馈可靠性", "伪阳性（flaky test）", "31.7% 漏检 / 73% 误拒", "验证器确定性输出"],
    ["修复指导", "仅知道失败", "依赖模型\"猜测\"", "后置条件/不变式/类型等精准诊断"],
]

# 表头
x_positions = [Inches(0.8), Inches(3.8), Inches(6.5), Inches(9.2)]
col_widths = [Inches(3.0), Inches(2.7), Inches(2.7), Inches(2.7)]
header_colors = [DARK_BLUE, MID_BLUE, MID_BLUE, ACCENT_ORANGE]

for j, (header, x, w, hc) in enumerate(zip(headers, x_positions, col_widths, header_colors)):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, Inches(1.6), w, Inches(0.5)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = hc
    rect.line.fill.background()
    tf = rect.text_frame
    tf.paragraphs[0].text = header
    tf.paragraphs[0].font.name = TITLE_FONT if j == 0 else BODY_FONT
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# 数据行
for i, row in enumerate(rows):
    y = Inches(2.2) + i * Inches(0.55)
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    for j, (cell, x, w) in enumerate(zip(row, x_positions, col_widths)):
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.55)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg
        rect.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        rect.line.width = Pt(0.5)
        
        tf = rect.text_frame
        tf.paragraphs[0].text = cell
        tf.paragraphs[0].font.name = TITLE_FONT if j == 0 else BODY_FONT
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = (j == 0)
        tf.paragraphs[0].font.color.rgb = BLACK if j != 3 else ACCENT_ORANGE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# 底部要点
add_textbox(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(0.4),
            "▎ 关键思路", TITLE_FONT, 16, True, DARK_BLUE)
points = [
    "形式化规约 = 精确的 \"正确性标准\"（无歧义、可证明）",
    "验证器反馈 = 结构化诊断信号的可靠来源",
    "\"规约→生成→验证→修复\" 闭环：用形式化验证替代不可靠的自审查",
]
add_bullet_textbox(slide, Inches(0.8), Inches(5.3), Inches(11), Inches(1.8),
                   points, BODY_FONT, 14, BLACK, "▸", 1.4)

print("形式化方案优势 ✓")


# ================================================================
# Slide 6: 研究目标与内容
# ================================================================
slide = add_blank_slide()
add_page_title(slide, "2. 研究目标与内容", "构建规约引导→代码生成→验证修复的端到端系统")

# 总体目标
add_rect(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.7), DARK_BLUE)
add_textbox(slide, Inches(1.0), Inches(1.55), Inches(11.3), Inches(0.6),
            "研究目标：从自然语言描述出发，自动生成 Dafny 形式化规约，在规约约束下生成代码，\n"
            "利用 Dafny 验证器的结构化反馈驱动多轮自修复，最终产出可证明正确的代码",
            TITLE_FONT, 13, True, WHITE, PP_ALIGN.LEFT, 1.3)

# 四个研究内容
contents = [
    ("内容一", "形式化规约自动生成", "NL→Dafny规约\n自校验+迭代修正", 
     "如何确保规约\n既完整又正确？"),
    ("内容二", "规约感知的代码生成", "规约约束融入prompt\n多候选生成+验证过滤",
     "规约表达方式对\n生成质量的影响？"),
    ("内容三", "验证反馈结构化解析", "错误类型分类\n错误定位+反例解释",
     "如何将验证器输出\n转化为修复引导？"),
    ("内容四", "多Agent退化感知修复", "5个Agent协同\nCoordinator控退化",
     "如何避免修复退化\n并最大化收敛概率？"),
]

for i, (num, title, approach, question) in enumerate(contents):
    x = Inches(0.6) + i * Inches(3.15)
    y = Inches(2.6)
    w = Inches(2.95)
    
    # 卡片
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(4.2)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_GRAY
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card.line.width = Pt(1)
    
    # 编号
    num_box = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x + w/2 - Inches(0.35), y + Inches(0.15), Inches(0.7), Inches(0.7)
    )
    colors = [DARK_BLUE, MID_BLUE, LIGHT_BLUE, ACCENT_ORANGE]
    num_box.fill.solid()
    num_box.fill.fore_color.rgb = colors[i]
    num_box.line.fill.background()
    tf = num_box.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.name = TITLE_FONT
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 标题
    add_textbox(slide, x + Inches(0.15), y + Inches(1.0), w - Inches(0.3), Inches(0.5),
                title, TITLE_FONT, 14, True, DARK_BLUE, PP_ALIGN.CENTER)
    
    add_rect(slide, x + Inches(0.5), y + Inches(1.55), w - Inches(1.0), Inches(0.02), ACCENT_ORANGE)
    
    # 方法
    add_textbox(slide, x + Inches(0.15), y + Inches(1.7), w - Inches(0.3), Inches(1.0),
                approach, BODY_FONT, 11, False, BLACK, PP_ALIGN.CENTER, 1.4)
    
    # 关键问题
    add_rect(slide, x + Inches(0.15), y + Inches(2.8), w - Inches(0.3), Inches(1.2), WHITE)
    add_textbox(slide, x + Inches(0.3), y + Inches(2.85), w - Inches(0.6), Inches(0.3),
                "❓ 关键科学问题", TITLE_FONT, 10, True, RED, PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.3), y + Inches(3.2), w - Inches(0.6), Inches(0.7),
                question, BODY_FONT, 11, False, GRAY, PP_ALIGN.CENTER, 1.3)

print("研究内容 ✓")


# ================================================================
# Slide 7: 系统架构 — 总览
# ================================================================
slide = add_blank_slide()
add_page_title(slide, "3. 技术路线与系统架构", "基于 LangGraph 的多 Agent 协同框架")

# 流程图
y_flow = Inches(1.6)
box_h = Inches(0.55)
box_w = Inches(1.8)
arrow_w = Inches(0.4)

# Phase I 标识
add_rect(slide, Inches(0.5), Inches(1.5), Inches(4.0), Inches(0.35), MID_BLUE)
add_textbox(slide, Inches(0.5), Inches(1.5), Inches(4.0), Inches(0.35),
            "Phase I: 规约生成与验证", TITLE_FONT, 11, True, WHITE, PP_ALIGN.CENTER)

# Phase II 标识
add_rect(slide, Inches(4.7), Inches(1.5), Inches(3.5), Inches(0.35), LIGHT_BLUE)
add_textbox(slide, Inches(4.7), Inches(1.5), Inches(3.5), Inches(0.35),
            "Phase II: 规约感知代码生成", TITLE_FONT, 11, True, WHITE, PP_ALIGN.CENTER)

# Phase III 标识
add_rect(slide, Inches(8.5), Inches(1.5), Inches(4.5), Inches(0.35), ACCENT_ORANGE)
add_textbox(slide, Inches(8.5), Inches(1.5), Inches(4.5), Inches(0.35),
            "Phase III: 验证反馈驱动修复（循环）", TITLE_FONT, 11, True, WHITE, PP_ALIGN.CENTER)

# Input
inp = add_flow_box(slide, Inches(0.5), Inches(2.2), Inches(2.0), box_h,
                    "输入：自然语言\n问题描述", MID_BLUE, WHITE, 11)

# Spec Agent
add_arrow_shape(slide, Inches(2.6), Inches(2.3), Inches(0.5), Inches(0.3), MID_BLUE)
spec = add_flow_box(slide, Inches(3.2), Inches(2.2), Inches(1.8), box_h,
                     "Spec Agent\nNL→Dafny规约", ACCENT_ORANGE, WHITE, 10)

# 规约验证
add_arrow_shape(slide, Inches(5.1), Inches(2.3), Inches(0.5), Inches(0.3), MID_BLUE)
spec_verify = add_flow_box(slide, Inches(5.7), Inches(2.2), Inches(1.8), box_h,
                            "Dafny\n规约校验", MID_BLUE, WHITE, 11)

# Code Agent
add_arrow_shape(slide, Inches(7.6), Inches(2.3), Inches(0.5), Inches(0.3), MID_BLUE)
code = add_flow_box(slide, Inches(8.2), Inches(2.2), Inches(1.8), box_h,
                     "Code Agent\n规约感知生成", ACCENT_ORANGE, WHITE, 10)

# Dafny Verify
add_arrow_shape(slide, Inches(10.1), Inches(2.3), Inches(0.5), Inches(0.3), MID_BLUE)
verify = add_flow_box(slide, Inches(10.7), Inches(2.2), Inches(1.8), box_h,
                       "Dafny\nVerify", RGBColor(0x1B, 0x78, 0x3D), WHITE, 11)

# 通过箭头
add_textbox(slide, Inches(10.7), Inches(2.85), Inches(1.8), Inches(0.3),
            "通过 → END ✓", BODY_FONT, 10, True, GREEN, PP_ALIGN.CENTER)

# 修复循环（下方）
y_repair = Inches(3.5)

add_rect(slide, Inches(5.2), y_repair - Inches(0.05), Inches(7.8), Inches(0.02), RGBColor(0xCC, 0x33, 0x33))

# 失败箭头说明
add_textbox(slide, Inches(7.8), y_repair - Inches(0.3), Inches(3.0), Inches(0.3),
            "❌ 失败 → 进入修复循环", BODY_FONT, 11, True, RED, PP_ALIGN.RIGHT)

# 修复子流程
diagnose = add_flow_box(slide, Inches(5.2), y_repair, Inches(1.8), box_h,
                         "Diagnose\nAgent (分析)", MID_BLUE, WHITE, 10)
add_arrow_shape(slide, Inches(7.1), y_repair + Inches(0.1), Inches(0.5), Inches(0.3), MID_BLUE)
repair = add_flow_box(slide, Inches(7.7), y_repair, Inches(1.8), box_h,
                       "Repair Agent\n(修复)", ACCENT_ORANGE, WHITE, 10)
add_arrow_shape(slide, Inches(9.6), y_repair + Inches(0.1), Inches(0.5), Inches(0.3), MID_BLUE)
coord = add_flow_box(slide, Inches(10.2), y_repair, Inches(1.8), box_h,
                      "Coordinator\n(退化检测)", DARK_BLUE, WHITE, 10)

# 循环回验证
add_rect(slide, Inches(11.9), y_repair + Inches(0.15), Inches(0.02), Inches(0.8), GRAY)
add_rect(slide, Inches(10.0), Inches(4.35), Inches(1.92), Inches(0.02), GRAY)
add_arrow_shape(slide, Inches(10.0), Inches(4.2), Inches(0.5), Inches(0.3), GRAY)
add_textbox(slide, Inches(10.5), Inches(4.15), Inches(1.5), Inches(0.3),
            "循环", BODY_FONT, 10, True, GRAY, PP_ALIGN.CENTER)

# Memory Agent
memory = add_flow_box(slide, Inches(0.5), y_repair + Inches(0.8), Inches(2.0), box_h,
                       "Memory Agent\n历史经验存储", MID_BLUE, WHITE, 10)
add_rect(slide, Inches(2.5), y_repair + Inches(1.05), Inches(0.02), Inches(0.5), GRAY)
add_arrow_shape(slide, Inches(2.5), y_repair + Inches(0.9), Inches(0.5), Inches(0.3), GRAY)
add_textbox(slide, Inches(3.0), y_repair + Inches(0.85), Inches(2.0), Inches(0.3),
            "反馈修复历史", BODY_FONT, 9, False, GRAY, PP_ALIGN.LEFT)

# 技术栈说明
add_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.4), LIGHT_GRAY)
add_textbox(slide, Inches(0.8), Inches(5.02), Inches(12), Inches(0.35),
            "技术栈：LangGraph（Agent编排）  |  DeepSeek / GPT（LLM）  |  Dafny CLI（验证器）  |  Python 3.11+",
            BODY_FONT, 12, False, GRAY, PP_ALIGN.CENTER)

# Agent职责说明
y_roles = Inches(5.6)
roles = [
    ("Spec Agent", "Dafny规约生成+自校验"),
    ("Code Agent", "规约约束下生成实现"),
    ("Dafny Verifier", "形式化验证"),
    ("Diagnose Agent", "结构化错误解析"),
    ("Repair Agent", "反馈驱动修复"),
    ("Coordinator", "退化检测+策略路由"),
]
for i, (name, duty) in enumerate(roles):
    x = Inches(0.5) + (i % 6) * Inches(2.1)
    add_textbox(slide, x, y_roles, Inches(2.0), Inches(0.25),
                f"● {name}：{duty}", BODY_FONT, 10, False, DARK_BLUE)

print("系统架构 ✓")


# ================================================================
# Slide 8: 关键技术 — 验证反馈结构化解析
# ================================================================
slide = add_blank_slide()
add_page_title(slide, "4. 关键技术：验证反馈的结构化解析", "将Dafny原始输出转化为LLM可用的修复指令")

# 左侧：原始输出
add_rect(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.5), LIGHT_GRAY)
add_textbox(slide, Inches(1.0), Inches(1.6), Inches(5.0), Inches(0.3),
            "🔴 Dafny 原始输出（非结构化）", TITLE_FONT, 14, True, RED)
raw_text = """dafny verify test.dfy

test.dfy(10,4): Error: A postcondition
  might not hold on this return path.
test.dfy(8,4): Related location:
  this is the postcondition that
  might not hold.
test.dfy(12,7): Error: This loop
  invariant might not be preserved
  by the loop.

Dafny 2 errors
  [31mResolution/type errors detected[0m"""
add_textbox(slide, Inches(1.0), Inches(2.0), Inches(5.0), Inches(2.8),
            raw_text, 'Courier New', 10, False, GRAY, PP_ALIGN.LEFT, 1.1)

# 箭头
add_arrow_shape(slide, Inches(6.4), Inches(2.5), Inches(0.6), Inches(0.35), ACCENT_ORANGE)
add_textbox(slide, Inches(6.4), Inches(2.9), Inches(0.6), Inches(0.3),
            "解析", BODY_FONT, 9, True, ACCENT_ORANGE, PP_ALIGN.CENTER)

# 右侧：结构化输出
add_rect(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(3.5), RGBColor(0xE8, 0xF5, 0xE9))
add_textbox(slide, Inches(7.4), Inches(1.6), Inches(5.0), Inches(0.3),
            "🟢 结构化错误信息", TITLE_FONT, 14, True, GREEN)
structured_text = """[
  {
    "error_type": "postcondition",
    "location": {"line": 10, "col": 4},
    "message": "A postcondition might not
      hold on this return path.",
    "related_spec": "ensures result >= x"
  },
  {
    "error_type": "invariant",
    "location": {"line": 12, "col": 7},
    "message": "This loop invariant might
      not be preserved by the loop.",
    "related": "invariant 0 <= i <= |s|"
  }
]"""
add_textbox(slide, Inches(7.4), Inches(2.0), Inches(5.0), Inches(2.8),
            structured_text, 'Courier New', 10, False, RGBColor(0x2E, 0x7D, 0x32), PP_ALIGN.LEFT, 1.1)

# 错误类型分类
y_types = Inches(5.3)
add_textbox(slide, Inches(0.8), y_types, Inches(3), Inches(0.3),
            "▎ 6种错误类型分类", TITLE_FONT, 15, True, DARK_BLUE)

error_types = [
    ("postcondition", "#CC3333"),
    ("precondition", "#E86C00"),
    ("invariant", "#2C5F8A"),
    ("syntax", "#663300"),
    ("type", "#336633"),
    ("other", "#666666"),
]
for i, (etype, color_str) in enumerate(error_types):
    x = Inches(0.8) + (i % 3) * Inches(3.0)
    y = y_types + Inches(0.45) + (i // 3) * Inches(0.4)
    r, g, b = int(color_str[1:3], 16), int(color_str[3:5], 16), int(color_str[5:7], 16)
    
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x, y + Inches(0.05), Inches(0.15), Inches(0.15)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = RGBColor(r, g, b)
    dot.line.fill.background()
    
    add_textbox(slide, x + Inches(0.25), y, Inches(2.5), Inches(0.3),
                etype, 'Courier New', 11, True, BLACK)

print("结构化解析 ✓")


# ================================================================
# Slide 9: 实验方案
# ================================================================
slide = add_blank_slide()
add_page_title(slide, "5. 实验方案", "数据集、Baseline、评估指标、消融实验")

# 数据集
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(3), Inches(0.35),
            "▎ 数据集", TITLE_FONT, 16, True, DARK_BLUE)

datasets = [
    ("HumanEval", "164题", "主实验"),
    ("MBPP", "974题", "泛化性验证"),
    ("Proving the Coding Interview", "27题", "质量评估"),
    ("自定义Dafny验证集", "50题", "消融实验"),
]
for i, (name, size, usage) in enumerate(datasets):
    x = Inches(0.8) + (i % 4) * Inches(3.1)
    y = Inches(1.95)
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.9), Inches(0.65)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = LIGHT_GRAY
    rect.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    rect.line.width = Pt(0.5)
    
    add_textbox(slide, x + Inches(0.1), y + Inches(0.02), Inches(2.7), Inches(0.25),
                f"{name}", TITLE_FONT, 11, True, DARK_BLUE, PP_ALIGN.LEFT)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.28), Inches(2.7), Inches(0.3),
                f"{size}   用途：{usage}", BODY_FONT, 10, False, GRAY, PP_ALIGN.LEFT)

# Baseline
add_textbox(slide, Inches(0.8), Inches(2.9), Inches(3), Inches(0.35),
            "▎ Baseline", TITLE_FONT, 16, True, DARK_BLUE)

baselines = [
    ("Direct Gen", "无规约无修复首轮生成", "Verif@1底线"),
    ("Self-Debug", "测试执行反馈驱动修复", "验证信号对比"),
    ("Reflexion", "语言强化学习修复", "修复策略对比"),
    ("ReflexiCoder", "RL内化反思-修复", "内化vs外部验证"),
    ("Single Agent", "本系统单Agent版本", "多Agent消融"),
]
for i, (name, desc, dim) in enumerate(baselines):
    x = Inches(0.8) + i * Inches(2.5)
    y = Inches(3.35)
    w = Inches(2.35)
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.8)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = LIGHT_GRAY
    rect.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    rect.line.width = Pt(0.5)
    
    add_textbox(slide, x + Inches(0.1), y + Inches(0.05), w - Inches(0.2), Inches(0.25),
                f"● {name}", TITLE_FONT, 11, True, DARK_BLUE)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.3), w - Inches(0.2), Inches(0.45),
                f"{desc}\n对比维度：{dim}", BODY_FONT, 9, False, GRAY, PP_ALIGN.LEFT, 1.2)

# 评估指标
add_textbox(slide, Inches(0.8), Inches(4.5), Inches(3), Inches(0.35),
            "▎ 评估指标", TITLE_FONT, 16, True, DARK_BLUE)

metrics = [
    ("Verif@1", "首轮验证通过率"),
    ("Verif@k", "k轮内验证通过率"),
    ("SpecAcc", "规约正确率（人工评定）"),
    ("AvgRounds", "平均修复轮次"),
    ("DegradationRate", "修复退化占比"),
    ("TokenCost", "每问题平均token消耗"),
]
for i, (mname, mdesc) in enumerate(metrics):
    x = Inches(0.8) + (i % 3) * Inches(3.1)
    y = Inches(5.0) + (i // 3) * Inches(0.4)
    add_textbox(slide, x, y, Inches(3.0), Inches(0.35),
                f"• {mname}：{mdesc}", BODY_FONT, 12, False, BLACK)

# 消融实验
add_textbox(slide, Inches(0.8), Inches(5.8), Inches(3), Inches(0.35),
            "▎ 5组消融实验", TITLE_FONT, 16, True, DARK_BLUE)

ablations = [
    "A1 规约有无 | A2 反馈结构化 | A3 Agent分离 | A4 轮次上限 | A5 模型对比"
]
add_bullet_textbox(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.6),
                   ablations, TITLE_FONT, 12, DARK_BLUE, "▸", 1.3)

print("实验方案 ✓")


# ================================================================
# Slide 10: 创新点
# ================================================================
slide = add_blank_slide()
add_page_title(slide, "6. 创新点", "三点核心创新")

innovations = [
    ("创新点一", "规约-代码协同生成方法",
     "将形式化规约自动生成与代码生成统一为协同优化过程",
     "现有工作：规约生成与代码生成分离（SpecGen生成规约后，手动或独立LLM调用生成代码）\n"
     "本方案：规约生成自带Dafny resolve双重校验，规约信息作为第一优先级约束嵌入代码生成Prompt",
     DARK_BLUE),
    ("创新点二", "验证反馈结构化解析框架",
     "建立从验证器原始输出到LLM可用修复指令的结构化翻译机制",
     "现有工作：直接投喂原始错误文本（'Error: postcondition might not hold'）\n"
     "本方案：6种错误类型分类 + 精确行列定位 + 关联规约提取 → Repair Agent精确修复",
     MID_BLUE),
    ("创新点三", "退化感知的多Agent迭代修复协议",
     "提出包含Coordinator主动控制退化的多Agent修复框架",
     "现有工作：固定轮次修复（3轮），无退化检测\n"
     "本方案：Coordinator监控修复质量变化，检测重复错误，自动切换修复策略",
     ACCENT_ORANGE),
]

for i, (num, title, summary, detail, color) in enumerate(innovations):
    y = Inches(1.6) + i * Inches(1.85)
    
    # 左侧色块
    add_rect(slide, Inches(0.8), y, Inches(1.2), Inches(1.6), color)
    add_textbox(slide, Inches(0.8), y + Inches(0.4), Inches(1.2), Inches(0.5),
                num, TITLE_FONT, 16, True, WHITE, PP_ALIGN.CENTER)
    
    # 主区域
    add_rect(slide, Inches(2.0), y, Inches(10.5), Inches(1.6), LIGHT_GRAY)
    
    # 标题
    add_textbox(slide, Inches(2.3), y + Inches(0.1), Inches(10), Inches(0.35),
                title, TITLE_FONT, 16, True, color)
    
    # 一句话概括
    add_textbox(slide, Inches(2.3), y + Inches(0.5), Inches(10), Inches(0.3),
                f"💡 {summary}", TITLE_FONT, 12, True, DARK_BLUE)
    
    # 对比说明
    add_textbox(slide, Inches(2.3), y + Inches(0.85), Inches(10), Inches(0.7),
                detail, BODY_FONT, 11, False, BLACK, PP_ALIGN.LEFT, 1.2)

print("创新点 ✓")


# ================================================================
# Slide 11: 研究计划
# ================================================================
slide = add_blank_slide()
add_page_title(slide, "7. 研究计划与进度安排", "8个月计划")

phases = [
    ("第1-2月", "基础搭建", LIGHT_BLUE,
     ["Dafny 验证器集成", "LangGraph 多Agent框架", "HumanEval→Dafny数据转换", "30个案例端到端跑通"]),
    ("第3-4月", "模块实现与实验", MID_BLUE,
     ["Spec Agent prompt设计与自校验", "Code Agent 规约感知生成", "Diagnose Agent 反馈解析", "HumanEval完整集初步实验"]),
    ("第5-6月", "方法完善与消融", DARK_BLUE,
     ["Coordinator Agent 退化检测", "Memory Agent 经验积累", "5组消融实验", "MBPP泛化性验证", "论文初稿"]),
    ("第7-8月", "论文撰写与投稿", ACCENT_ORANGE,
     ["论文修改完善", "补充Reviewer实验", "开源代码整理", "投稿CCF-B类会议/期刊"]),
]

for i, (period, title, color, tasks) in enumerate(phases):
    x = Inches(0.6) + i * Inches(3.2)
    y = Inches(1.6)
    w = Inches(3.0)
    
    # 头部
    add_rect(slide, x, y, w, Inches(0.6), color)
    add_textbox(slide, x, y + Inches(0.05), w, Inches(0.3),
                period, EN_FONT, 11, True, WHITE, PP_ALIGN.CENTER)
    add_textbox(slide, x, y + Inches(0.25), w, Inches(0.3),
                title, TITLE_FONT, 12, True, WHITE, PP_ALIGN.CENTER)
    
    # 内容
    add_rect(slide, x, y + Inches(0.6), w, Inches(3.0), LIGHT_GRAY)
    add_bullet_textbox(slide, x + Inches(0.15), y + Inches(0.7), w - Inches(0.3), Inches(2.8),
                       tasks, BODY_FONT, 11, BLACK, "▸", 1.4)

print("研究计划 ✓")


# ================================================================
# Slide 12: 实验验证 — 实际Demo展示
# ================================================================
slide = add_blank_slide()
add_page_title(slide, "8. 初步实验验证", "codegen-verify 原型系统已跑通完整闭环")

# 左侧：已完成工作
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.35),
            "▎ 已完成工作", TITLE_FONT, 16, True, GREEN)

done_items = [
    "基于 LangGraph 的多Agent框架完整搭建",
    "Dafny 验证器集成 + 结构化反馈解析（6种错误类型）",
    "Spec Agent：NL→Dafny规约 + 正则/Dafny resolve双重校验",
    "Code Agent：规约感知生成（含few-shot示例）",
    "Diagnose Agent：结构化错误诊断",
    "Repair Agent：反馈驱动自修复 + 重复错误检测",
    "Coordinator Agent：退化感知决策路由",
    "HumanEval 数据集转换与批量评测脚本",
]
add_bullet_textbox(slide, Inches(0.8), Inches(1.95), Inches(5.5), Inches(4.0),
                   done_items, BODY_FONT, 12, BLACK, "✓", 1.3)

# 右侧：技术亮点
add_textbox(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.35),
            "▎ 技术亮点", TITLE_FONT, 16, True, ACCENT_ORANGE)

highlights = [
    ("规约双重校验", "正则预检(var/for/while) + Dafny resolve"),
    ("嵌套循环桥接", "自动注入assert bridge连接内层循环结论"),
    ("退化检测降级", "连续重复错误自动切换修复策略"),
    ("历史累积", "Memory Agent记录修复经验"),
]

for i, (h_title, h_desc) in enumerate(highlights):
    y = Inches(1.95) + i * Inches(0.85)
    add_rect(slide, Inches(6.8), y, Inches(5.8), Inches(0.7), LIGHT_GRAY)
    add_textbox(slide, Inches(7.0), y + Inches(0.05), Inches(5.3), Inches(0.25),
                f"✦ {h_title}", TITLE_FONT, 12, True, ACCENT_ORANGE)
    add_textbox(slide, Inches(7.0), y + Inches(0.3), Inches(5.3), Inches(0.3),
                h_desc, BODY_FONT, 11, False, BLACK)

# 底部状态
add_rect(slide, Inches(0.8), Inches(6.5), Inches(11.8), Inches(0.6), DARK_BLUE)
add_textbox(slide, Inches(1.0), Inches(6.52), Inches(11.5), Inches(0.5),
            "🎯 已完成: 5个Agent端到端Pipeline全部跑通  |  正在进行: HumanEval批量评测与消融实验  |  后续: Baseline对比与论文撰写",
            TITLE_FONT, 13, True, WHITE, PP_ALIGN.CENTER)

print("初步验证 ✓")


# ================================================================
# Slide 13: 谢谢
# ================================================================
slide = add_blank_slide()
add_bg_rect(slide, DARK_BLUE)

add_rect(slide, 0, Inches(0.15), W, Inches(0.06), ACCENT_ORANGE)

add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
            "谢谢！请各位老师批评指正 🙏",
            TITLE_FONT, 40, True, WHITE, PP_ALIGN.CENTER)

add_rect(slide, Inches(5.5), Inches(3.3), Inches(2.3), Inches(0.04), ACCENT_ORANGE)

info_lines = [
    ("答辩人：Tez", BODY_FONT, 18, False, RGBColor(0xCC, 0xD5, 0xE8), PP_ALIGN.CENTER),
    ("", BODY_FONT, 8, False, WHITE, PP_ALIGN.CENTER),
    ("E-mail: tez@stu.ecnu.edu.cn", EN_FONT, 14, False, RGBColor(0x8A, 0xB4, 0xF8), PP_ALIGN.CENTER),
]
add_multiline_textbox(slide, Inches(3), Inches(3.8), Inches(7), Inches(1.5), info_lines)

add_rect(slide, 0, H - Inches(0.15), W, Inches(0.06), ACCENT_ORANGE)

print("谢谢 ✓")


# ========== 保存 ==========
output_path = "D:\\codegen-verify\\开题答辩PPT.pptx"
prs.save(output_path)
print(f"\n✅ PPT已保存到: {output_path}")
print(f"共 {len(prs.slides)} 页")

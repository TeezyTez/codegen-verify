"""
开题答辩PPT (v3 — 加大字体，填充排版)
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MID_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
LIGHT_BLUE = RGBColor(0x3A, 0x7C, 0xBD)
ACCENT_ORANGE = RGBColor(0xE8, 0x6C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
RED = RGBColor(0xCC, 0x33, 0x33)
GREEN = RGBColor(0x2E, 0x7D, 0x32)

TITLE_FONT = 'SimHei'
BODY_FONT = 'SimSun'
EN_FONT = 'Calibri'

# ====== Font sizes (all bumped up) ======
FS_PAGE_TITLE = 32
FS_SUBTITLE = 18
FS_SECTION = 22
FS_BODY_L = 18
FS_BODY_M = 16
FS_BODY_S = 14
FS_TABLE = 14
FS_FLOW = 14
FS_SMALL = 12

def blank():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    return s

def bg(s, c=DARK_BLUE):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()

def rect(s, l, t, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    return sh

def tb(s, l, t, w, h, txt, fn=BODY_FONT, fs=FS_BODY_M, b=False, c=BLACK, al=PP_ALIGN.LEFT, ls=1.3):
    bx = s.shapes.add_textbox(l, t, w, h)
    bx.text_frame.word_wrap = True
    p = bx.text_frame.paragraphs[0]
    p.text = txt; p.font.name = fn; p.font.size = Pt(fs); p.font.bold = b; p.font.color.rgb = c; p.alignment = al
    pPr = p._p.get_or_add_pPr()
    ln = pPr.makeelement(qn('a:lnSpc'), {})
    sp = ln.makeelement(qn('a:spcPct'), {'val': str(int(ls * 100000))})
    ln.append(sp); pPr.append(ln)
    return bx

def multi_tb(s, l, t, w, h, lines, ls=1.3):
    bx = s.shapes.add_textbox(l, t, w, h)
    bx.text_frame.word_wrap = True
    tf = bx.text_frame
    for i, (txt, fn, fs, b, c, al) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.font.name = fn; p.font.size = Pt(fs); p.font.bold = b; p.font.color.rgb = c; p.alignment = al
    return bx

def bullet(s, l, t, w, h, items, fn=BODY_FONT, fs=FS_BODY_M, c=BLACK, ch="●"):
    bx = s.shapes.add_textbox(l, t, w, h)
    bx.text_frame.word_wrap = True
    for i, item in enumerate(items):
        p = bx.text_frame.paragraphs[0] if i == 0 else bx.text_frame.add_paragraph()
        p.text = f"{ch} {item}"; p.font.name = fn; p.font.size = Pt(fs); p.font.color.rgb = c
    return bx

def flow(s, l, t, w, h, txt, bg_c=MID_BLUE, fg=WHITE, fs=FS_FLOW, fn=TITLE_FONT):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = bg_c; sh.line.fill.background()
    tf = sh.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = txt; tf.paragraphs[0].font.name = fn
    tf.paragraphs[0].font.size = Pt(fs); tf.paragraphs[0].font.color.rgb = fg
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return sh

def title_bar(s, txt, sub=None):
    rect(s, 0, 0, W, Inches(0.08), DARK_BLUE)
    tb(s, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7), txt, TITLE_FONT, FS_PAGE_TITLE, True, DARK_BLUE)
    if sub:
        tb(s, Inches(0.8), Inches(0.85), Inches(11), Inches(0.4), sub, BODY_FONT, FS_SUBTITLE, False, GRAY)
    rect(s, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.04), ACCENT_ORANGE)

# ================================================================
# S1: 封面
# ================================================================
s = blank(); bg(s, DARK_BLUE)
rect(s, 0, Inches(0.15), W, Inches(0.06), ACCENT_ORANGE)
tb(s, Inches(1), Inches(0.8), Inches(11), Inches(0.5), "华东师范大学", TITLE_FONT, 22, False, RGBColor(0x8A,0xB4,0xF8), PP_ALIGN.CENTER)
tb(s, Inches(1), Inches(1.8), Inches(11), Inches(1.2), "基于形式化契约引导的 LLM 代码生成\n与可验证自修复方法研究", TITLE_FONT, 38, True, WHITE, PP_ALIGN.CENTER, 1.4)
rect(s, Inches(5.5), Inches(3.3), Inches(2.3), Inches(0.04), ACCENT_ORANGE)
multi_tb(s, Inches(3), Inches(3.6), Inches(7), Inches(2.5), [
    ("硕士论文开题答辩", BODY_FONT, 20, False, RGBColor(0xCC,0xD5,0xE8), PP_ALIGN.CENTER),
    ("", BODY_FONT, 12, False, WHITE, PP_ALIGN.CENTER),
    ("答辩人：Tez", BODY_FONT, 18, False, WHITE, PP_ALIGN.CENTER),
    ("专  业：软件工程", BODY_FONT, 18, False, WHITE, PP_ALIGN.CENTER),
    ("导  师：XXX 教授", BODY_FONT, 18, False, WHITE, PP_ALIGN.CENTER),
])
rect(s, 0, H - Inches(0.15), W, Inches(0.06), ACCENT_ORANGE)
tb(s, Inches(1), H - Inches(0.6), Inches(11), Inches(0.4), "2026年6月  ·  华东师范大学", BODY_FONT, 16, False, RGBColor(0x8A,0xB4,0xF8), PP_ALIGN.CENTER)
print("1 封面 ✓")

# ================================================================
# S2: 目录
# ================================================================
s = blank(); title_bar(s, "目  录")
items = [("01", "研究背景与问题"), ("02", "高pass@1≠可证明正确性"), ("03", "研究目标与内容"),
         ("04", "技术路线与系统架构"), ("05", "实验方案与当前进展"), ("06", "创新点与研究计划")]
xs, ys, xg, yg = Inches(1.8), Inches(2.0), Inches(3.2), Inches(2.0)
for i, (num, t) in enumerate(items):
    col, row = i % 3, i // 3
    x, y = xs + col * xg, ys + row * yg
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.8), Inches(0.8))
    c.fill.solid()
    c.fill.fore_color.rgb = DARK_BLUE if row == 0 else MID_BLUE
    c.line.fill.background()
    tf = c.text_frame
    tf.paragraphs[0].text = num; tf.paragraphs[0].font.name = EN_FONT
    tf.paragraphs[0].font.size = Pt(20); tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tb(s, x + Inches(1.0), y + Inches(0.15), Inches(2.2), Inches(0.5), t, TITLE_FONT, FS_SECTION, True, DARK_BLUE)
print("2 目录 ✓")

# ================================================================
# S3: 研究背景
# ================================================================
s = blank(); title_bar(s, "1.1 研究背景", "LLM代码生成能力的飞跃，但正确性保障严重滞后")

tb(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5), "▎LLM代码生成能力的飞速提升", TITLE_FONT, FS_SECTION, True, DARK_BLUE)
bullet(s, Inches(0.8), Inches(2.1), Inches(5.5), Inches(2.5), [
    "Codex (2021): HumanEval Pass@1 = 28.8%",
    "Code Llama 70B (2023): Pass@1 = 67%",
    "ReflexiCoder (2026): Pass@1 = 94.51%",
    "GitHub Copilot: 超100万组织采用",
], fn=BODY_FONT, fs=FS_BODY_L)

tb(s, Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.5), "▎但正确性保障存在根本性鸿沟", TITLE_FONT, FS_SECTION, True, RED)
bullet(s, Inches(7.0), Inches(2.1), Inches(5.8), Inches(2.5), [
    "LLM 能 '生成' 代码，但本质上不 '理解' 执行语义",
    "生成代码可能包含隐藏的逻辑错误、边界条件遗漏",
    "Yu et al.: 测试用例从7.8扩至774个 → pass@1从85.4%骤降至33.5%",
    "关注点从 '能不能生成' 转向 '如何保证正确'",
], fn=BODY_FONT, fs=FS_BODY_L)

rect(s, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.0), LIGHT_GRAY)
tb(s, Inches(1.2), Inches(4.95), Inches(11), Inches(0.5), "💡 核心问题", TITLE_FONT, FS_SECTION, True, RED)
tb(s, Inches(1.2), Inches(5.5), Inches(11), Inches(1.0),
   "当代码生成能力接近天花板，我们能否保证生成的代码不仅仅是'看起来对'，而是'数学意义上正确'？\n"
   "形式化验证提供了答案：如果代码满足形式化规约，其正确性有数学证明。",
   BODY_FONT, FS_BODY_L, False, BLACK, PP_ALIGN.LEFT, 1.4)
print("3 背景 ✓")

# ================================================================
# S4: 三个可靠性危机
# ================================================================
s = blank(); title_bar(s, "1.2 三个可靠性危机", "现有自修复范式的根本性缺陷")
crises = [
    ("危机一：自审查双向失效", "🔍",
     ["Reddy et al. (2026): LLM 自审查漏检率达 31.7%",
      "Jin et al. (2026): 正确代码误判率飙升至 73%",
      "→ 模型在 '漏放' 与 '误拒' 上同时失效"]),
    ("危机二：反馈质量瓶颈", "🧪",
     ["Olausson et al. (2023): 人类反馈优于GPT-4自反馈1.58倍",
      "Arimbur et al. (2026): 断言错误修复率仅 45%",
      "→ 反馈质量而非生成能力才是自修复的核心瓶颈"]),
    ("危机三：缺乏形式化保障", "🛡️",
     ["测试执行 → 覆盖不全，非通过路径不可见",
      "LLM内省 → 共享生成者的认知盲区",
      "→ 缺少能提供数学级正确性保证的验证手段"]),
]
for i, (t, icon, items) in enumerate(crises):
    x = Inches(0.6) + i * Inches(4.2)
    w = Inches(3.9)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), w, Inches(4.8))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT_GRAY
    card.line.color.rgb = RGBColor(0xDD,0xDD,0xDD); card.line.width = Pt(1)
    tb(s, x + Inches(0.3), Inches(1.7), w - Inches(0.6), Inches(0.5), f"{icon} {t}", TITLE_FONT, FS_BODY_L, True, DARK_BLUE, PP_ALIGN.CENTER)
    rect(s, x + Inches(0.5), Inches(2.3), w - Inches(1.0), Inches(0.03), ACCENT_ORANGE)
    bullet(s, x + Inches(0.3), Inches(2.6), w - Inches(0.6), Inches(3.5), items, BODY_FONT, FS_BODY_M, BLACK, "▸")

rect(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.6), DARK_BLUE)
tb(s, Inches(0.8), Inches(6.65), Inches(11.7), Inches(0.5),
   "共性问题：自修复依赖的反馈信号本身不可靠 → 需要引入形式化验证提供可证明的正确性保证",
   TITLE_FONT, FS_BODY_M, True, WHITE, PP_ALIGN.CENTER)
print("4 三个危机 ✓")

# ================================================================
# S5: 高pass@1讨论
# ================================================================
s = blank(); title_bar(s, "1.3 高pass@1 ≠ 可证明的正确性", "为什么在 HumanEval 94% 的时代仍然需要本研究？")

args = [
    ("论点一：测试覆盖的 '假象'", GREEN,
     "HumanEval 每题平均仅 7.8 个测试用例。扩展到 774 个后，GPT-4 从 85.4% 骤降至 33.5%。\n高分在很大程度上是测试覆盖有限的假象。"),
    ("论点二：那 ~5-10% 的失败代码是关键", ACCENT_ORANGE,
     "当 LLM 犯错时：断言错误修复率仅 45%（Arimbur 2026），自审查漏检率 31.7%（Reddy 2026）。\n高 pass@1 并未解决这些失败代码的可靠修复问题。"),
    ("论点三：高 pass@1 是使能条件，不是替代方案", MID_BLUE,
     "正因为 LLM 足够强（94%），我们才有信心让它自动写规约。如果仍在 30% 水平，\nLLM 生成的规约也不可靠。本研究利用高 pass@1，叠加形式化保证。"),
]
for i, (t, c, d) in enumerate(args):
    y = Inches(1.6) + i * Inches(1.7)
    rect(s, Inches(0.8), y, Inches(0.08), Inches(1.4), c)
    tb(s, Inches(1.1), y + Inches(0.05), Inches(5.8), Inches(0.4), t, TITLE_FONT, FS_SECTION, True, c)
    tb(s, Inches(1.1), y + Inches(0.5), Inches(5.8), Inches(0.9), d, BODY_FONT, FS_BODY_M, False, BLACK, ls=1.3)

rect(s, Inches(7.5), Inches(1.5), Inches(5.2), Inches(5.5), LIGHT_GRAY)
tb(s, Inches(7.8), Inches(1.7), Inches(4.6), Inches(0.5), "🎯 一句话总结", TITLE_FONT, FS_SECTION, True, RED, PP_ALIGN.CENTER)
rect(s, Inches(8.2), Inches(2.3), Inches(3.8), Inches(0.03), ACCENT_ORANGE)
tb(s, Inches(7.8), Inches(2.6), Inches(4.6), Inches(2.5),
   "高 pass@1 回答的是\n'模型平均能写对多少'\n\n"
   "形式化验证回答的是\n'这段具体代码是否确实正确'\n\n"
   "二者正交，安全关键场景下后者不可或缺",
   TITLE_FONT, FS_BODY_L, True, DARK_BLUE, PP_ALIGN.CENTER, 1.5)
rect(s, Inches(7.8), Inches(5.5), Inches(4.6), Inches(1.2), WHITE)
tb(s, Inches(8.0), Inches(5.6), Inches(4.2), Inches(1.0),
   "本研究的定位：不是在 '平均能力' 上和 LLM 竞争，\n而是在 '单次验证' 中提供 LLM 无法提供的可证明正确性。",
   BODY_FONT, FS_BODY_S, False, GRAY, PP_ALIGN.CENTER, 1.3)
print("5 高pass@1讨论 ✓")

# ================================================================
# S6: 形式化方案优势
# ================================================================
s = blank(); title_bar(s, "1.4 为什么选择形式化方案？", "形式化规约提供了精确、无歧义的代码正确性标准")

headers = ["对比维度", "测试执行驱动", "LLM自审查驱动", "形式化验证（本方案）"]
rows = [
    ["正确性标准", "测试用例覆盖", "LLM 内省判断", "数学证明（全覆盖）"],
    ["反馈精度", "\"不通过\"（黑盒）", "\"可能有问题\"（模糊）", "结构化定位（行/列/类型）"],
    ["可靠性", "伪阳性（flaky）", "31.7%漏检/73%误拒", "验证器确定性输出"],
    ["修复指导", "仅知道失败", "依赖模型猜测", "后置条件/不变式等诊断"],
]
xps = [Inches(0.8), Inches(3.8), Inches(6.5), Inches(9.2)]
cws = [Inches(3.0), Inches(2.7), Inches(2.7), Inches(2.7)]
for j, (h, x, w, hc) in enumerate(zip(headers, xps, cws, [DARK_BLUE, MID_BLUE, MID_BLUE, ACCENT_ORANGE])):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.5), w, Inches(0.55))
    r.fill.solid(); r.fill.fore_color.rgb = hc; r.line.fill.background()
    tf = r.text_frame; tf.paragraphs[0].text = h
    tf.paragraphs[0].font.name = TITLE_FONT if j == 0 else BODY_FONT
    tf.paragraphs[0].font.size = Pt(FS_TABLE); tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

for i, row in enumerate(rows):
    y = Inches(2.15) + i * Inches(0.65)
    bg_c = LIGHT_GRAY if i % 2 == 0 else WHITE
    for j, (cell, x, w) in enumerate(zip(row, xps, cws)):
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.65))
        r.fill.solid(); r.fill.fore_color.rgb = bg_c
        r.line.color.rgb = RGBColor(0xDD,0xDD,0xDD); r.line.width = Pt(0.5)
        tf = r.text_frame; tf.paragraphs[0].text = cell
        tf.paragraphs[0].font.name = TITLE_FONT if j == 0 else BODY_FONT
        tf.paragraphs[0].font.size = Pt(FS_TABLE); tf.paragraphs[0].font.bold = (j == 0)
        tf.paragraphs[0].font.color.rgb = BLACK if j != 3 else ACCENT_ORANGE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

tb(s, Inches(0.8), Inches(4.8), Inches(11), Inches(0.5), "▎ 关键思路", TITLE_FONT, FS_SECTION, True, DARK_BLUE)
bullet(s, Inches(0.8), Inches(5.4), Inches(11), Inches(1.5), [
    "形式化规约 = 精确的 '正确性标准'（无歧义、可证明）",
    "验证器反馈 = 结构化诊断信号的可靠来源",
    "'规约→生成→验证→修复' 闭环：用形式化验证替代不可靠的自审查",
], BODY_FONT, FS_BODY_L, BLACK, "▸")
print("6 形式化方案 ✓")

# ================================================================
# S7: 研究目标与内容
# ================================================================
s = blank(); title_bar(s, "2. 研究目标与内容", "构建规约引导→代码生成→验证修复的端到端系统")
rect(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.7), DARK_BLUE)
tb(s, Inches(1.0), Inches(1.52), Inches(11.3), Inches(0.65),
   "研究目标：从自然语言描述出发→自动生成 Dafny 形式化规约→规约约束下生成代码→\n"
   "Dafny 验证器结构化反馈→驱动多轮自修复→产出可证明正确的代码",
   TITLE_FONT, FS_BODY_S, True, WHITE, PP_ALIGN.LEFT, 1.3)

contents = [
    ("内容一", "形式化规约自动生成", "NL→Dafny规约\n自校验+迭代修正", "如何确保规约\n既完整又正确？"),
    ("内容二", "规约感知的代码生成", "规约约束融入prompt\n多候选+验证过滤", "规约表达方式对\n生成质量的影响？"),
    ("内容三", "验证反馈结构化解析", "错误类型分类\n定位+反例解释", "如何将验证器输出\n转化为修复引导？"),
    ("内容四", "多Agent退化感知修复", "5个Agent协同\nCoordinator控退化", "如何避免修复退化\n最大化收敛概率？"),
]
colors = [DARK_BLUE, MID_BLUE, LIGHT_BLUE, ACCENT_ORANGE]
for i, (num, t, app, q) in enumerate(contents):
    x = Inches(0.5) + i * Inches(3.2)
    y = Inches(2.5)
    w = Inches(3.0)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(4.2))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT_GRAY
    card.line.color.rgb = RGBColor(0xDD,0xDD,0xDD); card.line.width = Pt(1)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x + w/2 - Inches(0.35), y + Inches(0.15), Inches(0.7), Inches(0.7))
    c.fill.solid(); c.fill.fore_color.rgb = colors[i]; c.line.fill.background()
    tf = c.text_frame; tf.paragraphs[0].text = num
    tf.paragraphs[0].font.name = TITLE_FONT; tf.paragraphs[0].font.size = Pt(FS_BODY_S)
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tb(s, x + Inches(0.15), y + Inches(1.0), w - Inches(0.3), Inches(0.5), t, TITLE_FONT, FS_BODY_M, True, DARK_BLUE, PP_ALIGN.CENTER)
    rect(s, x + Inches(0.5), y + Inches(1.55), w - Inches(1.0), Inches(0.03), ACCENT_ORANGE)
    tb(s, x + Inches(0.15), y + Inches(1.7), w - Inches(0.3), Inches(1.0), app, BODY_FONT, FS_BODY_S, False, BLACK, PP_ALIGN.CENTER, 1.3)
    rect(s, x + Inches(0.15), y + Inches(2.8), w - Inches(0.3), Inches(1.2), WHITE)
    tb(s, x + Inches(0.3), y + Inches(2.85), w - Inches(0.6), Inches(0.3), "❓ 关键科学问题", TITLE_FONT, FS_SMALL, True, RED, PP_ALIGN.CENTER)
    tb(s, x + Inches(0.3), y + Inches(3.2), w - Inches(0.6), Inches(0.7), q, BODY_FONT, FS_BODY_S, False, GRAY, PP_ALIGN.CENTER, 1.3)
print("7 研究内容 ✓")

# ================================================================
# S8: 系统架构
# ================================================================
s = blank(); title_bar(s, "3. 技术路线与系统架构", "基于 LangGraph 的多 Agent 协同框架")
for lbl, x, w, cl in [("Phase I: 规约生成与验证", Inches(0.5), Inches(4.0), MID_BLUE),
                       ("Phase II: 规约感知代码生成", Inches(4.7), Inches(3.5), LIGHT_BLUE),
                       ("Phase III: 验证反馈驱动修复（循环）", Inches(8.5), Inches(4.5), ACCENT_ORANGE)]:
    rect(s, x, Inches(1.5), w, Inches(0.4), cl)
    tb(s, x, Inches(1.5), w, Inches(0.4), lbl, TITLE_FONT, FS_SMALL, True, WHITE, PP_ALIGN.CENTER)

bh = Inches(0.6)
flow(s, Inches(0.5), Inches(2.2), Inches(2.0), bh, "输入: 自然语言\n问题描述", MID_BLUE, WHITE, FS_BODY_S)
flow(s, Inches(3.2), Inches(2.2), Inches(1.8), bh, "Spec Agent\nNL→Dafny规约", ACCENT_ORANGE, WHITE, FS_BODY_S)
flow(s, Inches(5.7), Inches(2.2), Inches(1.8), bh, "Dafny\n规约校验", MID_BLUE, WHITE, FS_BODY_S)
flow(s, Inches(8.2), Inches(2.2), Inches(1.8), bh, "Code Agent\n规约感知生成", ACCENT_ORANGE, WHITE, FS_BODY_S)
flow(s, Inches(10.7), Inches(2.2), Inches(1.8), bh, "Dafny\nVerify", RGBColor(0x1B,0x78,0x3D), WHITE, FS_BODY_S)
tb(s, Inches(10.7), Inches(2.85), Inches(1.8), Inches(0.3), "通过→ END ✓", BODY_FONT, FS_BODY_S, True, GREEN, PP_ALIGN.CENTER)

yr = Inches(3.6)
rect(s, Inches(5.2), yr - Inches(0.05), Inches(7.8), Inches(0.02), RED)
tb(s, Inches(7.8), yr - Inches(0.35), Inches(3.0), Inches(0.3), "❌ 失败→进入修复循环", BODY_FONT, FS_BODY_S, True, RED, PP_ALIGN.RIGHT)
flow(s, Inches(5.2), yr, Inches(1.8), bh, "Diagnose\nAgent (分析)", MID_BLUE, WHITE, FS_BODY_S)
flow(s, Inches(7.7), yr, Inches(1.8), bh, "Repair Agent\n(修复)", ACCENT_ORANGE, WHITE, FS_BODY_S)
flow(s, Inches(10.2), yr, Inches(1.8), bh, "Coordinator\n(退化检测)", DARK_BLUE, WHITE, FS_BODY_S)
flow(s, Inches(0.5), yr + Inches(0.8), Inches(2.0), bh, "Memory Agent\n历史经验存储", MID_BLUE, WHITE, FS_BODY_S)
# cycle arrow
rect(s, Inches(11.9), yr + Inches(0.15), Inches(0.02), Inches(0.8), GRAY)
rect(s, Inches(10.0), Inches(4.4), Inches(1.92), Inches(0.02), GRAY)
tb(s, Inches(10.5), Inches(4.2), Inches(1.5), Inches(0.3), "循环", BODY_FONT, FS_BODY_S, True, GRAY, PP_ALIGN.CENTER)

rect(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.4), LIGHT_GRAY)
tb(s, Inches(0.8), Inches(5.22), Inches(12), Inches(0.35),
   "技术栈：LangGraph (Agent编排)  |  DeepSeek / GPT (LLM)  |  Dafny CLI (验证器)  |  Python 3.11+",
   BODY_FONT, FS_BODY_S, False, GRAY, PP_ALIGN.CENTER)

roles = [("Spec Agent", "Dafny规约生成+自校验"), ("Code Agent", "规约约束下生成实现"),
         ("Dafny Verifier", "形式化验证"), ("Diagnose Agent", "结构化错误解析"),
         ("Repair Agent", "反馈驱动修复"), ("Coordinator", "退化检测+策略路由")]
for i, (n, d) in enumerate(roles):
    x = Inches(0.5) + (i % 6) * Inches(2.1)
    tb(s, x, Inches(5.85), Inches(2.0), Inches(0.3), f"● {n}：{d}", BODY_FONT, FS_SMALL, False, DARK_BLUE)
print("8 系统架构 ✓")

# ================================================================
# S9: 验证反馈结构化解析
# ================================================================
s = blank(); title_bar(s, "4. 关键技术：验证反馈的结构化解析", "将Dafny原始输出转化为LLM可用的修复指令")
rect(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.5), LIGHT_GRAY)
tb(s, Inches(1.0), Inches(1.6), Inches(5.0), Inches(0.35), "🔴 Dafny 原始输出（非结构化）", TITLE_FONT, FS_BODY_L, True, RED)
tb(s, Inches(1.0), Inches(2.1), Inches(5.0), Inches(2.8),
   "dafny verify test.dfy\n\ntest.dfy(10,4): Error: A postcondition\n  might not hold on this return path.\ntest.dfy(8,4): Related location:\n  this is the postcondition that\n  might not hold.\n\nDafny 2 errors",
   'Courier New', FS_BODY_S, False, GRAY, PP_ALIGN.LEFT, 1.1)

rect(s, Inches(7.2), Inches(1.5), Inches(5.5), Inches(3.5), RGBColor(0xE8,0xF5,0xE9))
tb(s, Inches(7.4), Inches(1.6), Inches(5.0), Inches(0.35), "🟢 结构化错误信息", TITLE_FONT, FS_BODY_L, True, GREEN)
tb(s, Inches(7.4), Inches(2.1), Inches(5.0), Inches(2.8),
   '[\n  {\n    "error_type": "postcondition",\n    "location": {"line": 10, "col": 4},\n    "related_spec": "ensures result >= x"\n  },\n  {\n    "error_type": "invariant",\n    "location": {"line": 12, "col": 7}\n  }\n]',
   'Courier New', FS_BODY_S, False, RGBColor(0x2E,0x7D,0x32), PP_ALIGN.LEFT, 1.1)

# 箭头
arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.4), Inches(2.5), Inches(0.6), Inches(0.35))
arr.fill.solid(); arr.fill.fore_color.rgb = ACCENT_ORANGE; arr.line.fill.background()
tb(s, Inches(6.4), Inches(2.9), Inches(0.6), Inches(0.3), "解析", BODY_FONT, FS_SMALL, True, ACCENT_ORANGE, PP_ALIGN.CENTER)

layers = [("① 错误分类", "postcondition/precondition/invariant/syntax/type"),
          ("② 精确定位", "行列号 + 关联规约片段提取"),
          ("③ 修复策略映射", "不同类型→差异化修复引导")]
for i, (lt, ld) in enumerate(layers):
    x = Inches(0.8) + (i % 3) * Inches(3.8)
    tb(s, x, Inches(5.4), Inches(3.5), Inches(0.35), lt, TITLE_FONT, FS_BODY_L, True, ACCENT_ORANGE)
    tb(s, x, Inches(5.8), Inches(3.5), Inches(0.5), ld, BODY_FONT, FS_BODY_S, False, GRAY)
print("9 结构化解析 ✓")

# ================================================================
# S10: 实验方案
# ================================================================
s = blank(); title_bar(s, "5. 实验方案", "数据集、Baseline、评估指标、消融实验")
tb(s, Inches(0.8), Inches(1.4), Inches(3), Inches(0.45), "▎ 数据集", TITLE_FONT, FS_SECTION, True, DARK_BLUE)
datasets = [("HumanEval", "164题", "主实验"), ("MBPP", "974题", "泛化性验证"),
            ("Proving Interview", "27题", "质量评估"), ("自定义Dafny集", "50题", "消融实验")]
for i, (n, sz, u) in enumerate(datasets):
    x = Inches(0.8) + (i % 4) * Inches(3.1)
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.95), Inches(2.9), Inches(0.7))
    r.fill.solid(); r.fill.fore_color.rgb = LIGHT_GRAY
    r.line.color.rgb = RGBColor(0xDD,0xDD,0xDD); r.line.width = Pt(0.5)
    tb(s, x + Inches(0.1), Inches(1.97), Inches(2.7), Inches(0.3), n, TITLE_FONT, FS_BODY_S, True, DARK_BLUE)
    tb(s, x + Inches(0.1), Inches(2.28), Inches(2.7), Inches(0.35), f"{sz}  用途：{u}", BODY_FONT, FS_SMALL, False, GRAY)

tb(s, Inches(0.8), Inches(2.9), Inches(3), Inches(0.45), "▎ Baseline", TITLE_FONT, FS_SECTION, True, DARK_BLUE)
baselines = [("Direct Gen", "无规约无修复", "Verif@1底线"), ("Self-Debug", "执行反馈修复", "信号对比"),
             ("Reflexion", "语言强化学习", "策略对比"), ("ReflexiCoder", "RL内化修复", "内化vs外部"),
             ("Single Agent", "单Agent版本", "多Agent消融")]
for i, (n, sd, dm) in enumerate(baselines):
    x = Inches(0.8) + i * Inches(2.5)
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.4), Inches(2.35), Inches(0.85))
    r.fill.solid(); r.fill.fore_color.rgb = LIGHT_GRAY
    r.line.color.rgb = RGBColor(0xDD,0xDD,0xDD); r.line.width = Pt(0.5)
    tb(s, x + Inches(0.1), Inches(3.42), Inches(2.15), Inches(0.3), f"● {n}", TITLE_FONT, FS_BODY_S, True, DARK_BLUE)
    tb(s, x + Inches(0.1), Inches(3.72), Inches(2.15), Inches(0.5), f"{sd}\n对比：{dm}", BODY_FONT, FS_SMALL, False, GRAY, ls=1.2)

tb(s, Inches(0.8), Inches(4.5), Inches(3), Inches(0.45), "▎ 评估指标", TITLE_FONT, FS_SECTION, True, DARK_BLUE)
for i, (mn, md) in enumerate([("Verif@1", "首轮验证通过率"), ("Verif@k", "k轮内通过率"),
    ("SpecAcc", "规约正确率"), ("AvgRounds", "平均修复轮次"),
    ("DegradationRate", "修复退化占比"), ("TokenCost", "每问题token消耗")]):
    x = Inches(0.8) + (i % 3) * Inches(3.1)
    tb(s, x, Inches(5.0) + (i // 3) * Inches(0.45), Inches(3.0), Inches(0.4), f"• {mn}：{md}", BODY_FONT, FS_BODY_M, False, BLACK)

tb(s, Inches(0.8), Inches(5.9), Inches(3), Inches(0.45), "▎ 消融实验", TITLE_FONT, FS_SECTION, True, DARK_BLUE)
bullet(s, Inches(0.8), Inches(6.35), Inches(11.5), Inches(0.6),
       ["A1 规约有无 | A2 反馈结构化 | A3 Agent分离 | A4 轮次上限 | A5 模型对比"],
       TITLE_FONT, FS_BODY_M, DARK_BLUE, "▸")
print("10 实验方案 ✓")

# ================================================================
# S11: 当前进展
# ================================================================
s = blank(); title_bar(s, "6. 当前实现进展", "核心模块已实现，系统整体进度约 60%")
hdrs = ["模块", "状态", "关键能力", "待完善项"]
rows = [("Spec Agent", "✅ v1.0", "NL→Dafny规约；双重校验", "规约完整性评估"),
        ("Code Agent", "✅ v1.0", "规约感知生成；assert桥接", "多候选生成+过滤"),
        ("Dafny Verifier", "✅ v1.0", "resolve+verify双阶段", "—"),
        ("Diagnose Agent", "✅ v1.0", "错误分类+分组分析", "输出结构化度提升"),
        ("Repair Agent", "✅ v1.0", "重复错误检测+换策略", "局部vs全局策略"),
        ("Coordinator", "⚠️ 部分", "条件路由已就绪", "退化检测+自适应"),
        ("Memory Agent", "❌ 待实现", "—", "ChromaDB检索"),
        ("HumanEval评测", "✅ 初步", "前5题跑通", "扩展至164题"),
        ("多模型对比", "❌ 未进行", "—", "DeepSeek/GPT/Qwen3")]
col_x = [Inches(0.8), Inches(3.0), Inches(4.4), Inches(7.8)]
col_w = [Inches(2.2), Inches(1.4), Inches(3.4), Inches(2.5)]
for j, (h, x, w) in enumerate(zip(hdrs, col_x, col_w)):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.5), w, Inches(0.45))
    r.fill.solid(); r.fill.fore_color.rgb = DARK_BLUE; r.line.fill.background()
    tf = r.text_frame; tf.paragraphs[0].text = h
    tf.paragraphs[0].font.name = TITLE_FONT; tf.paragraphs[0].font.size = Pt(FS_TABLE)
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER; tf.vertical_anchor = MSO_ANCHOR.MIDDLE

sc = {"✅": GREEN, "⚠️": ACCENT_ORANGE, "❌": RED}
for i, row in enumerate(rows):
    y = Inches(2.0) + i * Inches(0.5)
    bg_c = LIGHT_GRAY if i % 2 == 0 else WHITE
    for j, (cell, x, w) in enumerate(zip(row, col_x, col_w)):
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.5))
        r.fill.solid(); r.fill.fore_color.rgb = bg_c
        r.line.color.rgb = RGBColor(0xDD,0xDD,0xDD); r.line.width = Pt(0.5)
        tf = r.text_frame; tf.paragraphs[0].text = cell
        tf.paragraphs[0].font.name = TITLE_FONT if j == 0 else BODY_FONT
        tf.paragraphs[0].font.size = Pt(FS_SMALL); tf.paragraphs[0].font.bold = (j == 0)
        clr = BLACK
        if j == 1:
            for pk, pc in sc.items():
                if cell.startswith(pk): clr = pc; break
        tf.paragraphs[0].font.color.rgb = clr
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

rect(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.8), LIGHT_GRAY)
tb(s, Inches(1.0), Inches(6.22), Inches(11.3), Inches(0.75),
   "📊 基线实验：HumanEval 前5题，端到端通过率 40%（2/5），平均修复轮次 2.2\n"
   "简单问题均 1 轮通过，验证了 Pipeline 端到端收敛的有效性",
   BODY_FONT, FS_BODY_S, False, DARK_BLUE, PP_ALIGN.CENTER, 1.4)
print("11 当前进展 ✓")

# ================================================================
# S12: 创新点
# ================================================================
s = blank(); title_bar(s, "7. 创新点", "三点核心创新")
innovations = [
    ("创新点一", "规约-代码协同生成方法",
     "将形式化规约自动生成与代码生成统一为协同优化过程",
     "现有：规约生成与代码生成分离\n本方案：规约自带Dafny双重校验，作为第一优先级约束嵌入代码生成",
     DARK_BLUE),
    ("创新点二", "验证反馈结构化解析框架",
     "建立从验证器原始输出到LLM可用修复指令的翻译机制",
     "现有：直接投喂原始错误文本\n本方案：6种错误类型分类+精确定位+差异化修复策略映射",
     MID_BLUE),
    ("创新点三", "退化感知的多Agent修复协议",
     "提出包含Coordinator主动控制退化的多Agent修复框架",
     "现有：固定轮次修复，无退化检测\n本方案：Coordinator监控质量，检测重复错误，自动切换策略",
     ACCENT_ORANGE),
]
for i, (num, t, summary, detail, cl) in enumerate(innovations):
    y = Inches(1.5) + i * Inches(1.9)
    rect(s, Inches(0.8), y, Inches(1.2), Inches(1.7), cl)
    tb(s, Inches(0.8), y + Inches(0.5), Inches(1.2), Inches(0.5), num, TITLE_FONT, FS_BODY_L, True, WHITE, PP_ALIGN.CENTER)
    rect(s, Inches(2.0), y, Inches(10.5), Inches(1.7), LIGHT_GRAY)
    tb(s, Inches(2.3), y + Inches(0.1), Inches(10), Inches(0.4), t, TITLE_FONT, FS_SECTION, True, cl)
    tb(s, Inches(2.3), y + Inches(0.55), Inches(10), Inches(0.35), f"💡 {summary}", TITLE_FONT, FS_BODY_M, True, DARK_BLUE)
    tb(s, Inches(2.3), y + Inches(0.95), Inches(10), Inches(0.7), detail, BODY_FONT, FS_BODY_S, False, BLACK, PP_ALIGN.LEFT, 1.3)
print("12 创新点 ✓")

# ================================================================
# S13: 研究计划
# ================================================================
s = blank(); title_bar(s, "8. 研究计划与进度安排", "8个月计划")
phases = [
    ("第1-2月", "基础搭建", LIGHT_BLUE, ["Dafny验证器集成", "LangGraph多Agent框架", "HumanEval数据转换", "30个案例跑通"]),
    ("第3-4月", "模块实现+实验", MID_BLUE, ["Spec Agent Prompt设计", "Code Agent规约感知生成", "Diagnose Agent反馈解析", "HumanEval全量实验"]),
    ("第5-6月", "方法完善+消融", DARK_BLUE, ["Coordinator退化检测", "Memory Agent经验积累", "5组消融实验", "MBPP泛化性验证", "论文初稿"]),
    ("第7-8月", "论文撰写+投稿", ACCENT_ORANGE, ["论文修改完善", "补充Reviewer实验", "开源代码整理", "投稿CCF-B类会议"]),
]
for i, (period, t, cl, tasks) in enumerate(phases):
    x = Inches(0.5) + i * Inches(3.2)
    y = Inches(1.5)
    w = Inches(3.0)
    rect(s, x, y, w, Inches(0.65), cl)
    tb(s, x, y + Inches(0.03), w, Inches(0.3), period, EN_FONT, FS_BODY_S, True, WHITE, PP_ALIGN.CENTER)
    tb(s, x, y + Inches(0.28), w, Inches(0.3), t, TITLE_FONT, FS_BODY_S, True, WHITE, PP_ALIGN.CENTER)
    rect(s, x, y + Inches(0.65), w, Inches(3.0), LIGHT_GRAY)
    bullet(s, x + Inches(0.15), y + Inches(0.75), w - Inches(0.3), Inches(2.8), tasks, BODY_FONT, FS_BODY_S, BLACK, "▸")
print("13 研究计划 ✓")

# ================================================================
# S14: 谢谢
# ================================================================
s = blank(); bg(s, DARK_BLUE)
rect(s, 0, Inches(0.15), W, Inches(0.06), ACCENT_ORANGE)
tb(s, Inches(1), Inches(2.0), Inches(11), Inches(1.0), "谢谢！请各位老师批评指正 🙏", TITLE_FONT, 40, True, WHITE, PP_ALIGN.CENTER)
rect(s, Inches(5.5), Inches(3.3), Inches(2.3), Inches(0.04), ACCENT_ORANGE)
multi_tb(s, Inches(3), Inches(3.8), Inches(7), Inches(1.5), [
    ("答辩人：Tez", BODY_FONT, 20, False, RGBColor(0xCC,0xD5,0xE8), PP_ALIGN.CENTER),
    ("", BODY_FONT, 10, False, WHITE, PP_ALIGN.CENTER),
    ("E-mail: tez@stu.ecnu.edu.cn", EN_FONT, 16, False, RGBColor(0x8A,0xB4,0xF8), PP_ALIGN.CENTER),
])
rect(s, 0, H - Inches(0.15), W, Inches(0.06), ACCENT_ORANGE)
print("14 谢谢 ✓")

# ========== Save ==========
out = "D:\\codegen-verify\\开题答辩PPT_v3.pptx"
prs.save(out)
print(f"\n✅ PPT已保存: {out} 共 {len(prs.slides)} 页")

